from __future__ import annotations

import collections
import collections.abc

for name in ("Container", "Mapping", "MutableMapping", "Sequence", "MutableSequence", "Iterable", "Callable"):
    if not hasattr(collections, name) and hasattr(collections.abc, name):
        setattr(collections, name, getattr(collections.abc, name))

import logging
import os
import re
import requests
import uuid
from collections import OrderedDict
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any, Annotated, Dict, List, Literal, Optional, Union

from fastapi import APIRouter, HTTPException, Request
from fastapi.concurrency import run_in_threadpool
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from backend.chats.gemini_service import generate_response
from backend.chats.unsplash_service import fetch_unsplash_image, fetch_unsplash_url

logger = logging.getLogger(__name__)

APP_NAME = "Vitya Presentation API"
OUTPUT_DIR = Path(os.getenv("PPT_OUTPUT_DIR", "./outputs")).resolve()
DEFAULT_TEMPLATE_FILE = os.getenv("PPT_TEMPLATE_FILE", "./templates/base_template.pptx")
ASSET_DIR = Path(os.getenv("PPT_ASSET_DIR", "./assets")).resolve()
MAX_SLIDES = int(os.getenv("PPT_MAX_SLIDES", "30"))
MAX_BULLETS_PER_SLIDE = int(os.getenv("PPT_MAX_BULLETS_PER_SLIDE", "8"))
MAX_PARAGRAPH_CHARS = int(os.getenv("PPT_MAX_PARAGRAPH_CHARS", "900"))
ALLOW_ABSOLUTE_IMAGE_PATHS = os.getenv("PPT_ALLOW_ABSOLUTE_IMAGE_PATHS", "false").lower() == "true"

ALLOWED_SLIDE_TYPES = {
    "title_slide",
    "title_content",
    "mixed_content_slide",
    "chart_slide",
    "image_slide",
    "bullets_slide",
    "section_slide",
    "table_slide",
}

router = APIRouter()

_raw_cors = os.getenv("PPT_CORS_ORIGINS", "*")
CORS_ORIGINS = [x.strip() for x in _raw_cors.split(",") if x.strip()] or ["*"]

OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
ASSET_DIR.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------
# Schemas
# ---------------------------------------------------------------------

class SlidePluginText(BaseModel):
    type: Literal["text"]
    data: Dict[str, Any]


class SlidePluginBullets(BaseModel):
    type: Literal["bullets"]
    data: Dict[str, Any]


class SlidePluginParagraph(BaseModel):
    type: Literal["paragraph"]
    data: Dict[str, Any]


class SlidePluginChart(BaseModel):
    type: Literal["chart"]
    data: Dict[str, Any]


class SlidePluginImage(BaseModel):
    type: Literal["image"]
    data: Dict[str, Any]


class SlidePluginTable(BaseModel):
    type: Literal["table"]
    data: Dict[str, Any]


class SlidePluginNotes(BaseModel):
    type: Literal["notes"]
    data: Dict[str, Any]


class SlidePluginDiagram(BaseModel):
    type: Literal["diagram"]
    data: Dict[str, Any]


SlidePlugin = Annotated[
    Union[
        SlidePluginText,
        SlidePluginBullets,
        SlidePluginParagraph,
        SlidePluginChart,
        SlidePluginImage,
        SlidePluginTable,
        SlidePluginNotes,
        SlidePluginDiagram,
    ],
    Field(discriminator="type"),
]


class SlideSpec(BaseModel):
    layout: Optional[
        Literal[
            "title_slide",
            "title_content",
            "mixed_content_slide",
            "chart_slide",
            "image_slide",
            "bullets_slide",
            "section_slide",
            "table_slide",
        ]
    ] = None
    title: Optional[str] = None
    subtitle: Optional[str] = None
    title_color: Optional[str] = None
    title_font_size: Optional[int] = None
    title_bold: Optional[bool] = True
    title_align: Optional[str] = None
    subtitle_color: Optional[str] = None
    subtitle_font_size: Optional[int] = None
    subtitle_align: Optional[str] = None
    plugins: List[SlidePlugin] = Field(default_factory=list)


class PresentationPlan(BaseModel):
    title: str
    theme: Optional[Dict[str, str]] = None
    slides: List[SlideSpec]


class GenerateRequest(BaseModel):
    prompt: str = Field(..., min_length=1)
    template_name: Optional[str] = None
    slide_count: int = Field(default=8, ge=3, le=MAX_SLIDES)
    audience: Optional[str] = Field(default=None, max_length=120)
    tone: Optional[str] = Field(default=None, max_length=120)
    language: str = Field(default="English", max_length=80)
    include_citations: bool = False
    include_speaker_notes: bool = False
    use_gemini: bool = True

    include_title_slide: bool = True
    allow_bullets: bool = True
    allow_paragraph: bool = True
    allow_chart: bool = True
    allow_image: bool = True
    allow_section_slide: bool = True
    allow_table: bool = True

    background_theme: Optional[str] = None
    content_theme: Optional[str] = None
    visual_style: Optional[str] = None

    smart_mode: bool = True
    slide_types: Optional[List[str]] = None
    plan: Optional[PresentationPlan] = None


class GenerateResponse(BaseModel):
    job_id: str
    status: Literal["completed"]
    file_name: str
    download_url: str


class SaveResponse(BaseModel):
    presentation_id: str
    status: Literal["saved"]
    file_name: str
    download_url: str
    message: str


class RefineSlideRequest(BaseModel):
    text: str
    action: Optional[str] = "polish"


class RefineSlideResponse(BaseModel):
    refined_text: str



# ---------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------

def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def clean_ai_instructions(text: Optional[str]) -> str:
    if not text:
        return ""
    cleaned = normalize_whitespace(text)
    cleaned = re.sub(
        r"^(?i:\s*(?:break\s+down|explain\s+how|explain\s+the|explain|detail\s+the|detail|focus\s+on|highlight\s+the|highlight|describe\s+the|describe|conclude\s+with|discuss\s+the|discuss|provide\s+an|provide\s+a|provide|summarize\s+the|summarize)\s+)",
        "",
        cleaned,
    ).strip()
    return cleaned if cleaned else (text or "").strip()


def safe_filename(name: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9_-]+", "_", normalize_whitespace(name))[:80].strip("_")
    return cleaned or "presentation"


def normalize_slide_types(slide_types: Optional[List[str]]) -> Optional[List[str]]:
    if not slide_types:
        return None
    cleaned: List[str] = []
    for item in slide_types:
        value = normalize_whitespace(str(item)).lower()
        if value in ALLOWED_SLIDE_TYPES:
            cleaned.append(value)
    return cleaned or None


def resolve_template_path(template_name: Optional[str]) -> str:
    if template_name:
        candidate = Path(template_name).expanduser()
        if candidate.is_file():
            return str(candidate)
    return DEFAULT_TEMPLATE_FILE


def ensure_template_prs(template_file: str) -> Presentation:
    path = Path(template_file)
    prs = None
    if path.exists():
        try:
            prs = Presentation(str(path))
        except Exception as exc:
            logger.warning("Failed to load template %s: %s", path, exc)
    if prs is None:
        prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def ph(*names: str) -> tuple:
    values = []
    for name in names:
        value = getattr(PP_PLACEHOLDER, name, None)
        if value is not None:
            values.append(value)
    return tuple(values)


TITLE_PLACEHOLDER_TYPES = ph("TITLE", "CENTER_TITLE")
SUBTITLE_PLACEHOLDER_TYPES = ph("SUBTITLE")
BODY_PLACEHOLDER_TYPES = ph("BODY", "OBJECT", "VERTICAL_BODY", "VERTICAL_OBJECT", "TABLE")
IMAGE_PLACEHOLDER_TYPES = ph("PICTURE")
CHART_PLACEHOLDER_TYPES = ph("CHART")


def title_key(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", normalize_whitespace(text).lower()).strip()


def unique_title(title: str, suffix: str, seen: set[str]) -> str:
    base = normalize_whitespace(title) or "Slide"
    k = title_key(base)
    if k not in seen:
        seen.add(k)
        return base
    candidate = f"{base} - {suffix}"
    seen.add(title_key(candidate))
    return candidate


def split_text_into_chunks(text: str, max_chars: int = MAX_PARAGRAPH_CHARS) -> List[str]:
    text = normalize_whitespace(text)
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: List[str] = []
    current = ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if not current:
            current = sentence
        elif len(current) + 1 + len(sentence) <= max_chars:
            current += " " + sentence
        else:
            chunks.append(current)
            current = sentence
    if current:
        chunks.append(current)
    return chunks or [text[:max_chars]]


def chunk_list(items: List[Any], size: int) -> List[List[Any]]:
    size = max(1, size)
    return [items[i:i + size] for i in range(0, len(items), size)]


def parse_number(value: str) -> Optional[float]:
    try:
        return float(value)
    except Exception:
        return None


def hex_to_rgb(hex_color: str, default: Optional[RGBColor] = None) -> RGBColor:
    try:
        clean = (hex_color or "").replace("#", "").strip()
        if len(clean) == 3:
            clean = "".join(c * 2 for c in clean)
        return RGBColor.from_string(clean)
    except Exception:
        return default or RGBColor(15, 23, 42)


def is_light_color(rgb: RGBColor) -> bool:
    try:
        r, g, b = rgb[0], rgb[1], rgb[2]
        luminance = (0.299 * r + 0.587 * g + 0.114 * b)
        return luminance > 135
    except Exception:
        return False


def safe_list(value: Any) -> List[Any]:
    return value if isinstance(value, list) else []


@dataclass(frozen=True)
class Box:
    left: float
    top: float
    width: float
    height: float


def as_box(plan: Dict[str, Any], default: Box) -> Box:
    raw = plan.get("box") if isinstance(plan.get("box"), dict) else {}
    left = plan.get("left", raw.get("left", default.left))
    top = plan.get("top", raw.get("top", default.top))
    width = plan.get("width", raw.get("width", default.width))
    height = plan.get("height", raw.get("height", default.height))
    return Box(
        left=float(left),
        top=float(top),
        width=float(width),
        height=float(height),
    )


# ---------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------

THEME_COLORS = {
    "light": {"background": "F8FAFC", "gradient_start": "F8FAFC", "gradient_end": "E2E8F0", "accent": "2563EB", "text": "0F172A", "badge": "2563EB", "table_header_bg": "2563EB", "table_header_text": "FFFFFF", "table_row_bg1": "F8FAFC", "table_row_bg2": "E2E8F0", "table_row_text": "0F172A"},
    "dark": {"background": "0F172A", "gradient_start": "0F172A", "gradient_end": "31104B", "accent": "C084FC", "text": "FFFFFF", "badge": "C084FC", "table_header_bg": "31104B", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "0F172A", "table_row_text": "FFFFFF"},
    "midnight": {"background": "0F172A", "gradient_start": "0F172A", "gradient_end": "31104B", "accent": "C084FC", "text": "FFFFFF", "badge": "C084FC", "table_header_bg": "31104B", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "0F172A", "table_row_text": "FFFFFF"},
    "purple": {"background": "1E1B4B", "gradient_start": "1E1B4B", "gradient_end": "4C0519", "accent": "C084FC", "text": "FFFFFF", "badge": "C084FC", "table_header_bg": "4C0519", "table_header_text": "FFFFFF", "table_row_bg1": "2E1065", "table_row_bg2": "1E1B4B", "table_row_text": "FFFFFF"},
    "blue": {"background": "06101E", "gradient_start": "06101E", "gradient_end": "134074", "accent": "60A5FA", "text": "FFFFFF", "badge": "60A5FA", "table_header_bg": "134074", "table_header_text": "FFFFFF", "table_row_bg1": "0B2545", "table_row_bg2": "06101E", "table_row_text": "FFFFFF"},
    "ocean_blue": {"background": "06101E", "gradient_start": "06101E", "gradient_end": "134074", "accent": "38BDF8", "text": "FFFFFF", "badge": "38BDF8", "table_header_bg": "134074", "table_header_text": "FFFFFF", "table_row_bg1": "0B2545", "table_row_bg2": "06101E", "table_row_text": "FFFFFF"},
    "emerald": {"background": "022C22", "gradient_start": "022C22", "gradient_end": "047857", "accent": "34D399", "text": "FFFFFF", "badge": "34D399", "table_header_bg": "047857", "table_header_text": "FFFFFF", "table_row_bg1": "064E3B", "table_row_bg2": "022C22", "table_row_text": "FFFFFF"},
    "emerald_dark": {"background": "022C22", "gradient_start": "022C22", "gradient_end": "047857", "accent": "34D399", "text": "FFFFFF", "badge": "34D399", "table_header_bg": "047857", "table_header_text": "FFFFFF", "table_row_bg1": "064E3B", "table_row_bg2": "022C22", "table_row_text": "FFFFFF"},
    "cyberpunk_neon": {"background": "09090B", "gradient_start": "09090B", "gradient_end": "581C87", "accent": "F43F5E", "text": "FFFFFF", "badge": "F43F5E", "table_header_bg": "581C87", "table_header_text": "FFFFFF", "table_row_bg1": "2E1065", "table_row_bg2": "09090B", "table_row_text": "FFFFFF"},
    "wall_street": {"background": "022C22", "gradient_start": "022C22", "gradient_end": "1E293B", "accent": "10B981", "text": "FFFFFF", "badge": "10B981", "table_header_bg": "047857", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "022C22", "table_row_text": "FFFFFF"},
    "executive_gold": {"background": "1C1917", "gradient_start": "1C1917", "gradient_end": "78350F", "accent": "F59E0B", "text": "FFFFFF", "badge": "F59E0B", "table_header_bg": "78350F", "table_header_text": "FFFFFF", "table_row_bg1": "451A03", "table_row_bg2": "1C1917", "table_row_text": "FFFFFF"},
    "velvet_rose": {"background": "2A0813", "gradient_start": "2A0813", "gradient_end": "881337", "accent": "FB7185", "text": "FFFFFF", "badge": "FB7185", "table_header_bg": "881337", "table_header_text": "FFFFFF", "table_row_bg1": "4C0519", "table_row_bg2": "2A0813", "table_row_text": "FFFFFF"},
    "slate": {"background": "18181B", "gradient_start": "18181B", "gradient_end": "3F3F46", "accent": "A1A1AA", "text": "FFFFFF", "badge": "A1A1AA", "table_header_bg": "3F3F46", "table_header_text": "FFFFFF", "table_row_bg1": "27272A", "table_row_bg2": "18181B", "table_row_text": "FFFFFF"},
    "executive_slate": {"background": "18181B", "gradient_start": "18181B", "gradient_end": "3F3F46", "accent": "A1A1AA", "text": "FFFFFF", "badge": "A1A1AA", "table_header_bg": "3F3F46", "table_header_text": "FFFFFF", "table_row_bg1": "27272A", "table_row_bg2": "18181B", "table_row_text": "FFFFFF"},
    "titanium_white": {"background": "FFFFFF", "gradient_start": "FFFFFF", "gradient_end": "F4F4F5", "accent": "4F46E5", "text": "18181B", "badge": "4F46E5", "table_header_bg": "4F46E5", "table_header_text": "FFFFFF", "table_row_bg1": "F4F4F5", "table_row_bg2": "E4E4E7", "table_row_text": "18181B"},
    "sunset_glow": {"background": "2E1065", "gradient_start": "2E1065", "gradient_end": "9F1239", "accent": "FB7185", "text": "FFFFFF", "badge": "FB7185", "table_header_bg": "9F1239", "table_header_text": "FFFFFF", "table_row_bg1": "4C0519", "table_row_bg2": "2E1065", "table_row_text": "FFFFFF"},
    "ai": {"background": "0F172A", "gradient_start": "0F172A", "gradient_end": "31104B", "accent": "C084FC", "text": "F8FAFC", "badge": "C084FC", "table_header_bg": "31104B", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "0F172A", "table_row_text": "FFFFFF"},
    "data": {"background": "1E1B4B", "gradient_start": "1E1B4B", "gradient_end": "31104B", "accent": "C084FC", "text": "FFFFFF", "badge": "C084FC", "table_header_bg": "31104B", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "1E1B4B", "table_row_text": "FFFFFF"},
    "startup": {"background": "1E1B4B", "gradient_start": "1E1B4B", "gradient_end": "7C2D12", "accent": "F97316", "text": "FFFFFF", "badge": "F97316", "table_header_bg": "7C2D12", "table_header_text": "FFFFFF", "table_row_bg1": "431407", "table_row_bg2": "1E1B4B", "table_row_text": "FFFFFF"},
    "education": {"background": "FFFBEB", "gradient_start": "FFFBEB", "gradient_end": "FEF3C7", "accent": "D97706", "text": "451F00", "badge": "D97706", "table_header_bg": "D97706", "table_header_text": "FFFFFF", "table_row_bg1": "FEF3C7", "table_row_bg2": "FDE68A", "table_row_text": "451F00"},
    "finance": {"background": "0F172A", "gradient_start": "0F172A", "gradient_end": "14532D", "accent": "34D399", "text": "FFFFFF", "badge": "34D399", "table_header_bg": "14532D", "table_header_text": "FFFFFF", "table_row_bg1": "064E3B", "table_row_bg2": "0F172A", "table_row_text": "FFFFFF"},
    "medical": {"background": "FFF1F2", "gradient_start": "FFF1F2", "gradient_end": "FFE4E6", "accent": "E11D48", "text": "4C0519", "badge": "E11D48", "table_header_bg": "E11D48", "table_header_text": "FFFFFF", "table_row_bg1": "FFE4E6", "table_row_bg2": "FECDD3", "table_row_text": "4C0519"},
    "default": {"background": "0F172A", "gradient_start": "0F172A", "gradient_end": "31104B", "accent": "C084FC", "text": "FFFFFF", "badge": "C084FC", "table_header_bg": "31104B", "table_header_text": "FFFFFF", "table_row_bg1": "1E293B", "table_row_bg2": "0F172A", "table_row_text": "FFFFFF"},
}

THEME_KEYWORDS = {
    "ai": ["artificial intelligence", "machine learning", "deep learning", "neural", "llm", "genai", "generative ai", "model"],
    "data": ["data", "analytics", "dashboard", "sql", "etl", "visualization", "insight"],
    "startup": ["startup", "mvp", "founder", "pitch", "product launch", "scale"],
    "education": ["education", "school", "college", "student", "teacher", "course", "study"],
    "finance": ["finance", "money", "budget", "bank", "investment", "trading", "portfolio"],
    "medical": ["medical", "health", "doctor", "clinic", "hospital", "patient", "diagnosis"],
}

VISUAL_STYLES = {
    "minimal": {"show_top_bar": False, "shadow": False},
    "corporate": {"show_top_bar": True, "shadow": False},
    "academic": {"show_top_bar": False, "shadow": False},
    "modern_gradient": {"show_top_bar": True, "shadow": True},
}


def detect_theme(text: str) -> str:
    raw = normalize_whitespace(text or "").lower()
    normalized = re.sub(r"[^a-z0-9]+", " ", raw)
    padded = f" {normalized} "
    scores: Dict[str, int] = {}
    for theme, keywords in THEME_KEYWORDS.items():
        score = 0
        for kw in keywords:
            kw_norm = kw.lower().strip()
            if " " in kw_norm:
                if f" {kw_norm} " in padded:
                    score += 3
            else:
                if re.search(rf"\b{re.escape(kw_norm)}\b", normalized):
                    score += 1
        scores[theme] = score
    best_theme = max(scores, key=scores.get)
    return best_theme if scores[best_theme] > 0 else "default"


def detect_visual_style(text: str) -> str:
    raw = normalize_whitespace(text or "").lower()
    if any(word in raw for word in ("research", "study", "paper", "thesis", "seminar", "academic", "university")):
        return "academic"
    if any(word in raw for word in ("business", "company", "client", "report", "meeting", "corporate", "management")):
        return "corporate"
    if any(word in raw for word in ("modern", "ui", "design", "startup", "product", "demo", "landing", "brand")):
        return "modern_gradient"
    return "minimal"


def get_theme_palette(theme_input: Any) -> Dict[str, Any]:
    if isinstance(theme_input, dict):
        bg = theme_input.get("bg_color") or theme_input.get("background") or "#0F172A"
        g_start = theme_input.get("bg_gradient_start") or theme_input.get("bg_start") or bg
        g_end = theme_input.get("bg_gradient_end") or theme_input.get("bg_end") or bg
        txt = theme_input.get("text_color") or theme_input.get("text") or "#FFFFFF"
        acc = theme_input.get("accent_color") or theme_input.get("accent") or "#C084FC"
        bdg = theme_input.get("badge_color") or theme_input.get("slide_numbering_color") or theme_input.get("badge") or acc
        th_bg = theme_input.get("table_header_bg") or acc
        th_txt = theme_input.get("table_header_text") or "#FFFFFF"
        tr_bg1 = theme_input.get("table_row_bg1") or "#1E293B"
        tr_bg2 = theme_input.get("table_row_bg2") or bg
        tr_txt = theme_input.get("table_row_text") or txt
        return {
            "background": hex_to_rgb(bg),
            "gradient_start": hex_to_rgb(g_start),
            "gradient_end": hex_to_rgb(g_end),
            "text": hex_to_rgb(txt),
            "accent": hex_to_rgb(acc),
            "badge": hex_to_rgb(bdg),
            "table_header_bg": hex_to_rgb(th_bg),
            "table_header_text": hex_to_rgb(th_txt),
            "table_row_bg1": hex_to_rgb(tr_bg1),
            "table_row_bg2": hex_to_rgb(tr_bg2),
            "table_row_text": hex_to_rgb(tr_txt),
        }

    theme = normalize_whitespace(str(theme_input or "default")).lower()
    raw = THEME_COLORS.get(theme, THEME_COLORS["default"])
    bg_hex = raw["background"]
    g_start_hex = raw.get("gradient_start", bg_hex)
    g_end_hex = raw.get("gradient_end", bg_hex)
    badge_hex = raw.get("badge") or raw.get("accent")

    th_bg_hex = raw.get("table_header_bg") or raw.get("accent")
    th_txt_hex = raw.get("table_header_text") or "FFFFFF"
    tr_bg1_hex = raw.get("table_row_bg1") or "1E293B"
    tr_bg2_hex = raw.get("table_row_bg2") or bg_hex
    tr_txt_hex = raw.get("table_row_text") or raw.get("text")

    return {
        "background": hex_to_rgb(bg_hex),
        "gradient_start": hex_to_rgb(g_start_hex),
        "gradient_end": hex_to_rgb(g_end_hex),
        "accent": hex_to_rgb(raw["accent"]),
        "text": hex_to_rgb(raw["text"]),
        "badge": hex_to_rgb(badge_hex),
        "table_header_bg": hex_to_rgb(th_bg_hex),
        "table_header_text": hex_to_rgb(th_txt_hex),
        "table_row_bg1": hex_to_rgb(tr_bg1_hex),
        "table_row_bg2": hex_to_rgb(tr_bg2_hex),
        "table_row_text": hex_to_rgb(tr_txt_hex),
    }


def get_visual_style(style_name: Optional[str]) -> Dict[str, Any]:
    style = normalize_whitespace(style_name or "minimal").lower()
    return VISUAL_STYLES.get(style, VISUAL_STYLES["minimal"])


def apply_background_theme(slide, theme_input: Any, visual_style: Optional[str] = None) -> None:
    palette = get_theme_palette(theme_input)
    fill = slide.background.fill

    try:
        if palette.get("gradient_start") and palette.get("gradient_end") and palette["gradient_start"] != palette["gradient_end"]:
            fill.gradient()
            fill.gradient_angle = 135.0
            stops = fill.gradient_stops
            stops[0].position = 0.0
            stops[0].color.rgb = palette["gradient_start"]
            stops[1].position = 1.0
            stops[1].color.rgb = palette["gradient_end"]
            return
    except Exception:
        pass

    fill.solid()
    fill.fore_color.rgb = palette["background"]


def set_run_style(run, font_size: int, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def configure_text_frame(tf, *, font_size: int, color: Optional[RGBColor] = None, bold: bool = False) -> None:
    try:
        tf.word_wrap = True
    except Exception:
        pass
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    try:
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
    except Exception:
        pass

    for p in tf.paragraphs:
        for run in p.runs:
            set_run_style(run, font_size=font_size, bold=bold, color=color)


def best_font_size_for_bullets(points: List[Any], base: int = 18) -> int:
    count = max(1, len(points))
    longest = max((len(normalize_whitespace(str(p))) for p in points), default=0)
    size = base
    if count >= 10:
        size -= 2
    if count >= 7:
        size -= 1
    if longest >= 100:
        size -= 2
    elif longest >= 70:
        size -= 1
    return max(16, size)


def best_font_size_for_paragraph(text: str, base: int = 20) -> int:
    text = normalize_whitespace(text)
    size = base
    if len(text) > 400:
        size -= 3
    elif len(text) > 250:
        size -= 2
    elif len(text) > 150:
        size -= 1
    return max(16, size)


def find_placeholder_by_types(slide, placeholder_types: tuple) -> Optional[Any]:
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type in placeholder_types:
                return shape
        except Exception:
            continue
    return None


def set_shape_text(shape, text: str, font_size: int = 20, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_style(run, font_size=font_size, bold=bold, color=color)
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def add_textbox(slide, left, top, width, height, text: str, font_size: int = 20, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_style(run, font_size=font_size, bold=bold, color=color)
    try:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def write_text_or_fallback(
    slide,
    placeholder_types: tuple,
    text: str,
    *,
    fallback_left: float,
    fallback_top: float,
    fallback_width: float,
    fallback_height: float,
    font_size: int,
    bold: bool = False,
    color: Optional[RGBColor] = None,
) -> None:
    shape = find_placeholder_by_types(slide, placeholder_types)
    if shape is not None:
        try:
            set_shape_text(shape, text, font_size=font_size, bold=bold, color=color)
            return
        except Exception:
            pass

    add_textbox(
        slide,
        Inches(fallback_left),
        Inches(fallback_top),
        Inches(fallback_width),
        Inches(fallback_height),
        text,
        font_size=font_size,
        bold=bold,
        color=color,
    )


def set_slide_notes(slide, notes: str) -> None:
    notes = normalize_whitespace(notes)
    if not notes:
        return
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        tf.clear()
        tf.text = notes
    except Exception:
        pass


def sanitize_image_path(path_text: str) -> Optional[str]:
    path_text = normalize_whitespace(path_text)
    if not path_text:
        return None

    if path_text.startswith("http://") or path_text.startswith("https://"):
        url_clean = path_text.split("?")[0].rstrip("/")
        file_part = url_clean.split("/")[-1] or "web_image"
        safe_name = safe_filename(file_part)
        target_file = ASSET_DIR / f"download_{safe_name}.jpg"
        if target_file.exists() and target_file.stat().st_size > 1000:
            return str(target_file)
        try:
            resp = requests.get(path_text, timeout=10)
            if resp.status_code == 200 and len(resp.content) > 1000:
                target_file.write_bytes(resp.content)
                return str(target_file)
        except Exception as exc:
            logger.warning("Failed to download image from URL %s: %s", path_text, exc)

    candidate = Path(path_text)
    if candidate.is_absolute():
        if ALLOW_ABSOLUTE_IMAGE_PATHS and candidate.exists():
            return str(candidate)
        return None

    resolved = (ASSET_DIR / candidate).resolve()
    try:
        if resolved.exists():
            return str(resolved)
    except Exception:
        pass
    return None


def fetch_unsplash_image_for_topic(topic_query: str) -> Optional[str]:
    """Fetch an image from Unsplash service for a given slide topic or title."""
    return fetch_unsplash_image(topic_query)


# ---------------------------------------------------------------------
# Layout detection
# ---------------------------------------------------------------------

@dataclass
class LayoutSpec:
    layout_index: int


@dataclass
class LayoutInfo:
    index: int
    name: str
    placeholder_count: int
    has_title: bool
    has_body: bool
    has_picture: bool
    has_chart: bool
    has_table: bool


class TemplateDetector:
    def __init__(self, template_file: str):
        self.template_file = template_file
        self.prs = ensure_template_prs(template_file)
        self.layouts = self._inspect_layouts()

    def _inspect_layouts(self) -> List[LayoutInfo]:
        layouts: List[LayoutInfo] = []
        for idx, layout in enumerate(self.prs.slide_layouts):
            name = (layout.name or f"layout_{idx}").strip().lower()
            has_title = False
            has_body = False
            has_picture = False
            has_chart = False
            has_table = False

            for ph_shape in layout.placeholders:
                try:
                    ph_type = ph_shape.placeholder_format.type
                    if ph_type in TITLE_PLACEHOLDER_TYPES:
                        has_title = True
                    elif ph_type in SUBTITLE_PLACEHOLDER_TYPES or ph_type in BODY_PLACEHOLDER_TYPES:
                        has_body = True
                    elif ph_type in IMAGE_PLACEHOLDER_TYPES:
                        has_picture = True
                    elif ph_type in CHART_PLACEHOLDER_TYPES:
                        has_chart = True
                    elif ph_type == getattr(PP_PLACEHOLDER, "TABLE", None):
                        has_table = True
                except Exception:
                    continue

            layouts.append(
                LayoutInfo(
                    index=idx,
                    name=name,
                    placeholder_count=len(layout.placeholders),
                    has_title=has_title,
                    has_body=has_body,
                    has_picture=has_picture,
                    has_chart=has_chart,
                    has_table=has_table,
                )
            )
        return layouts

    def _score(self, layout: LayoutInfo, target: str) -> int:
        name = layout.name
        score = 0

        if target == "title_slide":
            if any(k in name for k in ("title", "cover", "front")):
                score += 100
            if layout.has_title:
                score += 30
            if not layout.has_body:
                score += 10

        elif target in {"title_content", "bullets_slide"}:
            if any(k in name for k in ("content", "body", "text", "bullet", "list")):
                score += 100
            if layout.has_title and layout.has_body:
                score += 40
            if layout.placeholder_count >= 2:
                score += 10

        elif target == "section_slide":
            if any(k in name for k in ("section", "divider", "break")):
                score += 100
            if layout.has_title and not layout.has_body:
                score += 20

        elif target == "chart_slide":
            if any(k in name for k in ("chart", "graph", "data", "analytics")):
                score += 100
            if layout.has_chart:
                score += 40
            if layout.has_title and layout.has_body:
                score += 15

        elif target == "image_slide":
            if any(k in name for k in ("image", "picture", "photo", "visual")):
                score += 100
            if layout.has_picture:
                score += 40
            if layout.has_title and layout.has_body:
                score += 15

        elif target == "table_slide":
            if any(k in name for k in ("table", "data", "content", "body")):
                score += 100
            if layout.has_table or layout.has_body:
                score += 25
            if layout.has_title:
                score += 10

        if layout.placeholder_count == 0:
            score -= 20

        return score

    def pick_best_layout_index(self, target: str) -> int:
        if not self.layouts:
            return 0
        ranked = sorted(
            ((self._score(layout, target), layout.index) for layout in self.layouts),
            key=lambda x: x[0],
            reverse=True,
        )
        best_score, best_index = ranked[0]
        return best_index if best_score > 0 else 0

    def build_registry(self) -> Dict[str, LayoutSpec]:
        return {
            "title_slide": LayoutSpec(self.pick_best_layout_index("title_slide")),
            "title_content": LayoutSpec(self.pick_best_layout_index("title_content")),
            "section_slide": LayoutSpec(self.pick_best_layout_index("section_slide")),
            "bullets_slide": LayoutSpec(self.pick_best_layout_index("bullets_slide")),
            "chart_slide": LayoutSpec(self.pick_best_layout_index("chart_slide")),
            "image_slide": LayoutSpec(self.pick_best_layout_index("image_slide")),
            "table_slide": LayoutSpec(self.pick_best_layout_index("table_slide")),
            "mixed_content_slide": LayoutSpec(self.pick_best_layout_index("title_content")),
        }


@lru_cache(maxsize=16)
def get_layout_registry(template_file: str) -> Dict[str, LayoutSpec]:
    return TemplateDetector(template_file).build_registry()


# ---------------------------------------------------------------------
# Mixed slide geometry
# ---------------------------------------------------------------------

class MixedLayoutResolver:
    """
    Intelligent layout resolver for mixed slide content.

    Slide content types supported:
        diagram, paragraph, bullets, chart, table, image

    Coordinates are based on a 16:9 presentation.
    """

    # Full usable content area
    FULL = Box(0.8, 1.4, 11.7, 5.3)

    # Global geometry
    LEFT = 0.8
    TOP = 1.4
    WIDTH = 11.7
    HEIGHT = 5.3
    GAP = 0.2

    # Relative vertical-space priority for dynamic layouts.
    # Higher value = more space.
    HEIGHT_WEIGHT = {
        "diagram": 1.0,
        "paragraph": 0.8,
        "bullets": 1.0,
        "chart": 1.3,
        "table": 1.4,
        "image": 1.2,
    }

    @staticmethod
    def resolve_list(plugin_types: List[str]) -> List[Box]:
        """
        Resolve layout boxes for an ordered list of plugin types.
        Handles duplicates (e.g. multiple paragraphs or bullets) without overlapping.
        """
        if not plugin_types:
            return []
        unique_kinds = set(plugin_types)
        if len(unique_kinds) == len(plugin_types):
            dict_res = MixedLayoutResolver.resolve(unique_kinds)
            return [dict_res.get(t, MixedLayoutResolver.FULL) for t in plugin_types]

        count = len(plugin_types)
        total_gap = MixedLayoutResolver.GAP * (count - 1)
        available_height = max(1.0, MixedLayoutResolver.HEIGHT - total_gap)
        weights = [MixedLayoutResolver.HEIGHT_WEIGHT.get(k, 1.0) for k in plugin_types]
        total_weight = sum(weights) or 1.0

        boxes: List[Box] = []
        current_top = MixedLayoutResolver.TOP
        for weight in weights:
            h = (available_height * weight) / total_weight
            boxes.append(
                Box(
                    MixedLayoutResolver.LEFT,
                    round(current_top, 2),
                    MixedLayoutResolver.WIDTH,
                    round(h, 2),
                )
            )
            current_top += h + MixedLayoutResolver.GAP
        return boxes

    @staticmethod
    def resolve(plugin_types: set[str]) -> Dict[str, Box]:
        """
        Resolve layout boxes for a set of content/plugin types.
        """

        kinds = set(plugin_types)

        if not kinds:
            return {}

        # -------------------------------------------------------------
        # Single content type
        # -------------------------------------------------------------
        if len(kinds) == 1:
            kind = next(iter(kinds))
            return {
                kind: MixedLayoutResolver.FULL
            }

        # -------------------------------------------------------------
        # Diagram + Bullets
        # -------------------------------------------------------------
        if kinds == {"diagram", "bullets"}:
            return {
                "diagram": Box(
                    0.8, 1.4, 11.7, 1.45
                ),
                "bullets": Box(
                    0.8, 3.05, 11.7, 3.65
                ),
            }

        # -------------------------------------------------------------
        # Diagram + Paragraph
        # -------------------------------------------------------------
        if kinds == {"diagram", "paragraph"}:
            return {
                "diagram": Box(
                    0.8, 1.4, 11.7, 1.45
                ),
                "paragraph": Box(
                    0.8, 3.05, 11.7, 3.65
                ),
            }

        # -------------------------------------------------------------
        # Diagram + Paragraph + Bullets
        # -------------------------------------------------------------
        if kinds == {
            "diagram",
            "paragraph",
            "bullets",
        }:
            return {
                "diagram": Box(
                    0.8, 1.4, 11.7, 1.30
                ),
                "paragraph": Box(
                    0.8, 2.85, 11.7, 1.30
                ),
                "bullets": Box(
                    0.8, 4.35, 11.7, 2.35
                ),
            }

        # -------------------------------------------------------------
        # Paragraph + Image
        # -------------------------------------------------------------
        if kinds == {"paragraph", "image"}:
            return {
                "paragraph": Box(
                    0.8, 1.5, 5.6, 4.8
                ),
                "image": Box(
                    6.7, 1.5, 5.8, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Paragraph + Bullets
        # -------------------------------------------------------------
        if kinds == {"paragraph", "bullets"}:
            return {
                "paragraph": Box(
                    0.8, 1.4, 11.7, 1.45
                ),
                "bullets": Box(
                    0.8, 3.05, 11.7, 3.65
                ),
            }

        # -------------------------------------------------------------
        # Image + Bullets
        # -------------------------------------------------------------
        if kinds == {"image", "bullets"}:
            return {
                "image": Box(
                    0.8, 1.5, 5.6, 4.8
                ),
                "bullets": Box(
                    6.7, 1.5, 5.8, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Paragraph + Chart
        # -------------------------------------------------------------
        if kinds == {"paragraph", "chart"}:
            return {
                "paragraph": Box(
                    0.8, 1.5, 5.0, 4.8
                ),
                "chart": Box(
                    6.1, 1.5, 6.4, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Bullets + Chart
        # -------------------------------------------------------------
        if kinds == {"bullets", "chart"}:
            return {
                "bullets": Box(
                    0.8, 1.5, 5.0, 4.8
                ),
                "chart": Box(
                    6.1, 1.5, 6.4, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Paragraph + Table
        # -------------------------------------------------------------
        if kinds == {"table", "paragraph"}:
            return {
                "paragraph": Box(
                    0.8, 1.4, 11.7, 1.25
                ),
                "table": Box(
                    0.8, 2.85, 11.7, 3.85
                ),
            }

        # -------------------------------------------------------------
        # Bullets + Table
        # -------------------------------------------------------------
        if kinds == {"table", "bullets"}:
            return {
                "bullets": Box(
                    0.8, 1.4, 11.7, 1.55
                ),
                "table": Box(
                    0.8, 3.15, 11.7, 3.55
                ),
            }

        # -------------------------------------------------------------
        # Image + Chart
        # -------------------------------------------------------------
        if kinds == {"image", "chart"}:
            return {
                "image": Box(
                    0.8, 1.5, 5.6, 4.8
                ),
                "chart": Box(
                    6.7, 1.5, 5.8, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Image + Table
        # -------------------------------------------------------------
        if kinds == {"image", "table"}:
            return {
                "image": Box(
                    0.8, 1.5, 5.3, 4.8
                ),
                "table": Box(
                    6.3, 1.5, 6.2, 4.8
                ),
            }

        # -------------------------------------------------------------
        # Chart + Table
        # -------------------------------------------------------------
        if kinds == {"chart", "table"}:
            return {
                "chart": Box(
                    0.8, 1.5, 5.3, 4.8
                ),
                "table": Box(
                    6.3, 1.5, 6.2, 4.8
                ),
            }

        # -------------------------------------------------------------
        # 3-column: image + chart + bullets
        # -------------------------------------------------------------
        if kinds == {
            "image",
            "chart",
            "bullets",
        }:
            return {
                "image": Box(
                    0.8, 1.5, 3.65, 4.8
                ),
                "chart": Box(
                    4.75, 1.5, 3.65, 4.8
                ),
                "bullets": Box(
                    8.70, 1.5, 3.8, 4.8
                ),
            }

        # -------------------------------------------------------------
        # 3-block vertical:
        # paragraph + bullets + table
        # -------------------------------------------------------------
        if kinds == {
            "paragraph",
            "bullets",
            "table",
        }:
            return {
                "paragraph": Box(
                    0.8, 1.4, 11.7, 1.15
                ),
                "bullets": Box(
                    0.8, 2.75, 11.7, 1.65
                ),
                "table": Box(
                    0.8, 4.60, 11.7, 2.10
                ),
            }

        # -------------------------------------------------------------
        # 3-block vertical:
        # paragraph + bullets + chart
        # -------------------------------------------------------------
        if kinds == {
            "paragraph",
            "bullets",
            "chart",
        }:
            return {
                "paragraph": Box(
                    0.8, 1.4, 11.7, 1.10
                ),
                "bullets": Box(
                    0.8, 2.70, 5.2, 3.95
                ),
                "chart": Box(
                    6.25, 2.70, 6.25, 3.95
                ),
            }

        # -------------------------------------------------------------
        # Dynamic fallback
        # -------------------------------------------------------------
        return MixedLayoutResolver._dynamic_layout(kinds)

    # -----------------------------------------------------------------
    # Dynamic fallback
    # -----------------------------------------------------------------

    @staticmethod
    def _dynamic_layout(kinds: set[str]) -> Dict[str, Box]:
        """
        Dynamic non-overlapping layout.

        Uses content-specific weights rather than giving every
        component the same height.
        """

        ordered_kinds = [
            k
            for k in [
                "diagram",
                "paragraph",
                "bullets",
                "chart",
                "table",
                "image",
            ]
            if k in kinds
        ]

        # Unknown plugin types
        if not ordered_kinds:
            ordered_kinds = list(kinds)

        count = len(ordered_kinds)

        if count == 1:
            return {
                ordered_kinds[0]: MixedLayoutResolver.FULL
            }

        # Total gap between components
        total_gap = MixedLayoutResolver.GAP * (count - 1)

        # Available height after gaps
        available_height = (
            MixedLayoutResolver.HEIGHT - total_gap
        )

        # Calculate weights
        weights = [
            MixedLayoutResolver.HEIGHT_WEIGHT.get(
                kind,
                1.0,
            )
            for kind in ordered_kinds
        ]

        total_weight = sum(weights)

        result: Dict[str, Box] = {}

        current_top = MixedLayoutResolver.TOP

        for kind, weight in zip(
            ordered_kinds,
            weights,
        ):
            height = (
                available_height
                * weight
                / total_weight
            )

            result[kind] = Box(
                MixedLayoutResolver.LEFT,
                round(current_top, 2),
                MixedLayoutResolver.WIDTH,
                round(height, 2),
            )

            current_top += height + MixedLayoutResolver.GAP

        return result


# ---------------------------------------------------------------------
# Planner
# ---------------------------------------------------------------------

def build_gemini_slide_script(req: GenerateRequest) -> Optional[str]:
    from dotenv import load_dotenv
    load_dotenv()
    if not req.use_gemini or not os.getenv("GEMINI_API_KEY"):
        return None

    options = []
    if req.audience:
        options.append(f"Target Audience: {req.audience}.")
    if req.tone:
        options.append(f"Tone: {req.tone}.")
    options.append(f"Language: {req.language}.")
    if req.include_citations:
        options.append("Add a final Sources slide containing short, credible source names/URLs. Do not invent fake citations.")
    if req.include_speaker_notes:
        options.append("Add one concise `Notes:` line to every non-title slide.")

    instructions = " ".join(options)
    gemini_prompt = f"""You are a world-class presentation designer and domain content strategist.
Create a professional, visually rich, logically structured PowerPoint script about this topic:
"{req.prompt}"

{instructions}

Follow these strict design and content rules:

1. TOPIC & STRUCTURE INTELLIGENCE:
- Automatically select the ideal slide count (target around {req.slide_count} slides unless topic warrants slightly more or fewer).
- Adapt the structure based on the topic domain:
  * Machine Learning / Tech: Title -> Problem vs Traditional -> Types/Architecture -> ML Pipeline (Diagram) -> Model Evaluation -> Real-world Applications -> Challenges -> Future Scope -> Conclusion.
  * Electric Vehicles / Product: Title -> Overview -> Powertrain Architecture -> Battery & Charging Tech -> EV vs ICE Comparison (Table) -> Market Adoption Trend (Chart) -> Infrastructure -> Future -> Conclusion.
  * History / Science: Title -> Background -> Timeline/Key Milestones -> Major Events -> Personalities/Components -> Impact/Significance -> Legacy/Conclusion.
  * Cyber Security: Title -> Threat Landscape -> Attack Vectors (Diagram) -> Defense Mechanisms -> Real-world Case Study -> Best Practices -> Future Trends -> Conclusion.
  * Business / Financial: Title -> Executive Summary -> Market Analysis -> Business Model -> Growth Metrics (Chart with Units) -> Competitor Matrix (Table) -> Financial Outlook -> Conclusion.

2. CONTENT QUALITY & DENSITY:
- No long paragraphs. Use 3-5 concise, informative bullet points per slide (max 15-20 words per bullet).
- Never place internal instructions like "Break down...", "Explain...", "Detail...", "Focus on...", "Highlight...", "Describe..." in titles or content. Provide audience-ready text ONLY.
- Explain technical terms clearly when introduced.
- Do NOT fabricate stats, dates, research findings, or fake citations. Mark illustrative trends with "[Illustrative Data]".

3. VISUAL MATCHING & DIAGRAMS:
- Match slide intent to the visual structure:
  * Process / Workflow / Architecture -> Visual Diagram (e.g. `Diagram: [Input Data] ➔ [Preprocessing] ➔ [Model Training] ➔ [Evaluation] ➔ [Deployment]`)
  * Feature Comparison -> Multi-criteria Comparison Table with headings (e.g. `Feature | Option A | Option B | Criterion`)
  * Data Trend -> Chart with metric name, axis units, and series name.
  * Key Visual Highlights / Innovations / Applications / Case Studies -> Include topic-relevant Unsplash keyword for `Image: [specific query]` on 2 to 4 slides across the presentation to make it visually engaging!

4. CHARTS & DATA INTEGRITY:
- NEVER output empty values or all-zero placeholders ("Phase 1: 0"). Every chart MUST specify realistic, domain-relevant numerical data and clear category names matching the topic!
- Choose ideal chart type: `Chart: column`, `Chart: line`, `Chart: bar`, `Chart: pie`, `Chart: area`, or `Chart: donut`.
- Examples:
  * Deep-Sea Submergence: `Series Name: Operational Depth (Meters)`, `1960 (Trieste): 10916`, `1989 (Shinkai): 6527`, `2012 (Challenger): 10908`, `2019 (Limiting Factor): 10928`
  * AI & ML Models: `Series Name: Parameter Scale (Billions)`, `2018 (GPT-1): 0.11`, `2019 (GPT-2): 1.5`, `2020 (GPT-3): 175`, `2023 (GPT-4): 1800`
  * Business & Tech: `Series Name: Market Growth ($ Millions)`, `2021: 15.4`, `2022: 28.6`, `2023: 45.2`, `2024: 68.9`

5. REAL COMPARISON TABLES:
- Comparison slides MUST feature an actual data table comparing options against criteria (e.g. Cost, Performance, Security, Architecture, Scalability).

OUTPUT FORMAT:
Return ONLY the plain-text slide script. Do not use Markdown code fences, introductory prose, or JSON.
Use these exact slide format structures:

Slide 1:
Title: [Specific Professional Title]
Subtitle: [Informative Subtitle]

Slide 2:
Title: [Topic Overview / Executive Summary]
Paragraph: [2-3 concise sentences summarizing core context.]

Slide 3:
Title: [Architecture / Workflow]
Diagram: [Input Component] ➔ [Processing Layer] ➔ [Core Engine] ➔ [Output/Result]
Bullets:
- Key architectural component 1
- Key architectural component 2
- Key architectural component 3

Slide 4:
Title: [Technology / Model Comparison]
Table:
Criterion | Solution A | Solution B | Solution C
Performance | High (99.9% uptime) | Medium (98.5%) | High (99.5%)
Cost | Enterprise Tier | Pay-as-you-go | Open Source
Scalability | Multi-region | Single-region | Hybrid Cloud

Slide 5:
Title: [Market Trend / Performance Metric]
Chart: column
Series Name: Market Adoption Rate (%) [Illustrative Data]
2021: 15
2022: 28
2023: 45
2024: 68
2025: 85

Slide 6:
Title: [Visual Innovation Showcase]
Image: [Specific topic keyword image query]
Paragraph: [Context explaining the visual innovation...]

Notes: [Concise speaker note for the presenter.]"""

    response = generate_response(gemini_prompt)
    if not response or response.startswith("Gemini API key is not configured") or response.startswith("Gemini error"):
        logger.warning("Gemini presentation planning failed; using local planner")
        return None
    returned_slides = len(re.findall(r"(?im)^\s*slide\s*\d+\s*[:\-]", response))
    if returned_slides < 2:
        logger.warning("Gemini returned an invalid presentation script; using local planner")
        return None
    if returned_slides != req.slide_count:
        logger.warning(
            "Gemini returned %s slides when %s were requested; using the valid response",
            returned_slides,
            req.slide_count,
        )
    return response

class PromptPlanner:
    def plan(
        self,
        prompt: str,
        *,
        include_title_slide: bool = True,
        allow_bullets: bool = True,
        allow_paragraph: bool = True,
        allow_chart: bool = True,
        allow_image: bool = True,
        allow_section_slide: bool = True,
        allow_table: bool = True,
        smart_mode: bool = True,
        slide_types: Optional[List[str]] = None,
        target_slide_count: Optional[int] = None,
    ) -> PresentationPlan:
        prompt = (prompt or "").strip()
        prompt_l = prompt.lower()
        blocks = self.extract_structured_slides(prompt)
        presentation_title = self.extract_overall_title(prompt, blocks)
        allowed_set = set(slide_types) if slide_types else None

        def allowed(layout_name: str) -> bool:
            if allowed_set is not None and layout_name not in allowed_set:
                return False
            if layout_name == "title_slide":
                return include_title_slide
            if layout_name == "bullets_slide":
                return allow_bullets
            if layout_name == "title_content":
                return allow_paragraph or allow_bullets or allow_image or allow_chart or allow_table
            if layout_name == "mixed_content_slide":
                return allow_paragraph or allow_bullets or allow_image or allow_chart or allow_table
            if layout_name == "chart_slide":
                return allow_chart
            if layout_name == "image_slide":
                return allow_image
            if layout_name == "section_slide":
                return allow_section_slide
            if layout_name == "table_slide":
                return allow_table
            return True

        slides: List[SlideSpec] = []

        if not blocks:
            if include_title_slide:
                slides.append(self._make_title_slide(presentation_title))

            if allow_section_slide and any(k in prompt_l for k in ("overview", "introduction", "intro")):
                slides.append(self._make_section_slide("Overview"))

            if allow_bullets:
                slides.append(self._make_bullets_slide(
                    "Key Points",
                    ["Clear problem statement", "Main ideas and workflow", "Practical use cases"],
                ))

            if allow_paragraph and len(prompt) > 80:
                slides.append(self._make_paragraph_slide("Summary", prompt[:600]))

            if allow_chart and re.search(r"\b(chart|graph|trend|growth|comparison)\b", prompt_l):
                slides.append(self._make_chart_slide(
                    "Trend Chart",
                    {
                        "chart_type": "line",
                        "title": "Trend Chart",
                        "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
                        "values": [12, 18, 24, 31, 39, 48],
                        "series_name": "Usage",
                    },
                ))

            # Keep non-AI fallback diverse with varied slide layouts
            fallback_topics = [
                ("Introduction", "paragraph"),
                ("Why It Matters", "bullets"),
                ("Key Features & Specifications", "table"),
                ("Growth & Impact Metrics", "chart"),
                ("Best Practices & Use Cases", "mixed"),
                ("Challenges & Solutions", "bullets"),
                ("Implementation Steps", "paragraph"),
                ("Conclusion & Next Steps", "bullets"),
            ]
            desired_count = min(max(target_slide_count or len(slides), 1), MAX_SLIDES)
            for idx, (topic, layout_type) in enumerate(fallback_topics):
                if len(slides) >= desired_count:
                    break
                if layout_type == "paragraph" and allow_paragraph:
                    slides.append(self._make_paragraph_slide(topic, f"Key context and strategic overview regarding {topic.lower()}."))
                elif layout_type == "table" and allow_table:
                    slides.append(self._make_table_slide(
                        topic,
                        {
                            "title": topic,
                            "headers": ["Metric / Phase", "Standard", "Target"],
                            "rows": [["Phase 1", "Initial Setup", "Completed"], ["Phase 2", "Optimization", "In Progress"], ["Phase 3", "Scaling", "Planned"]],
                        }
                    ))
                elif layout_type == "chart" and allow_chart:
                    slides.append(self._make_chart_slide(
                        topic,
                        {
                            "chart_type": "column",
                            "title": topic,
                            "categories": ["Phase 1", "Phase 2", "Phase 3", "Phase 4"],
                            "values": [25, 55, 80, 100],
                            "series_name": "Progress %",
                        }
                    ))
                elif layout_type == "mixed" and allow_paragraph and allow_bullets:
                    slides.append(self._make_mixed_slide(
                        topic,
                        [
                            SlidePluginParagraph(type="paragraph", data={"text": f"Overview of {topic.lower()} in modern workflow."}),
                            SlidePluginBullets(type="bullets", data={"points": ["Key advantage 1", "Key advantage 2", "Key advantage 3"]}),
                        ]
                    ))
                elif allow_bullets:
                    slides.append(self._make_bullets_slide(
                        topic,
                        [
                            f"How {topic.lower()} relates to the presentation topic",
                            "Important context and practical considerations",
                            "A concise takeaway for the audience",
                        ],
                    ))
                elif allow_paragraph:
                    slides.append(self._make_paragraph_slide(topic, f"Key context about {topic.lower()}."))

            return PresentationPlan(title=presentation_title, slides=slides[:MAX_SLIDES])

        seen_titles: set[str] = set()

        for idx, block in enumerate(blocks):
            parsed = self.parse_slide_block(block)

            raw_title = (
                parsed["title"]
                or self.extract_heading_title(block)
                or (presentation_title if idx == 0 else f"Slide {idx + 1}")
            )
            raw_title = normalize_whitespace(raw_title)

            notes = normalize_whitespace(parsed.get("notes", ""))
            paragraph = normalize_whitespace(parsed.get("paragraph", ""))
            bullets = parsed.get("bullets") or []
            image_path = normalize_whitespace(parsed.get("image_path") or "")
            chart_series = parsed.get("chart_series") or {}
            chart_points = parsed.get("chart_points") or []
            table_rows = parsed.get("table_rows") or []

            if idx == 0 and allowed("title_slide"):
                slides.append(self._make_title_slide(raw_title, parsed.get("subtitle", "")))

            if allow_section_slide and self.is_section_block(block, parsed):
                slides.append(self._make_section_slide(parsed.get("section_title") or raw_title))
                continue

            def t(suffix: str) -> str:
                return unique_title(raw_title, suffix, seen_titles)

            plugins: List[SlidePlugin] = []

            diagram = normalize_whitespace(parsed.get("diagram", ""))
            if diagram:
                plugins.append(SlidePluginDiagram(type="diagram", data={"diagram": diagram}))

            if paragraph and allow_paragraph:
                plugins.append(SlidePluginParagraph(type="paragraph", data={"text": paragraph, "font_size": 18}))

            if bullets and allow_bullets:
                plugins.append(SlidePluginBullets(type="bullets", data={"points": bullets[:MAX_BULLETS_PER_SLIDE]}))

            if image_path and allow_image:
                plugins.append(SlidePluginImage(type="image", data={"path": image_path, "caption": raw_title}))

            if (chart_series or chart_points) and allow_chart:
                plugins.append(SlidePluginChart(type="chart", data=self.build_chart_payload(parsed)))

            if table_rows and allow_table:
                plugins.append(SlidePluginTable(type="table", data=self.build_table_payload(parsed, title=raw_title)))

            if len(plugins) >= 2:
                box_list = MixedLayoutResolver.resolve_list([p.type for p in plugins])
                adjusted: List[SlidePlugin] = []

                for plugin, box in zip(plugins, box_list):
                    data = dict(plugin.data)
                    if box:
                        data["box"] = {"left": box.left, "top": box.top, "width": box.width, "height": box.height}
                    data.pop("title", None)
                    adjusted.append(type(plugin)(type=plugin.type, data=data))

                if notes:
                    adjusted.append(SlidePluginNotes(type="notes", data={"notes": notes}))

                slides.append(
                    SlideSpec(
                        layout="mixed_content_slide",
                        title=t("Overview"),
                        plugins=adjusted,
                    )
                )
                continue

            if len(plugins) == 1:
                plugin = plugins[0]
                if plugin.type == "paragraph":
                    slides.append(self._make_paragraph_slide(t("Overview"), paragraph, notes))
                elif plugin.type == "bullets":
                    slides.append(self._make_bullets_slide(t("Key Points"), bullets, notes))
                elif plugin.type == "image":
                    slides.append(self._make_image_slide(t("Visual"), image_path))
                elif plugin.type == "chart":
                    slides.append(self._make_chart_slide(t("Chart"), self.build_chart_payload(parsed)))
                elif plugin.type == "table":
                    slides.append(self._make_table_slide(t("Table"), self.build_table_payload(parsed, title=raw_title)))
                continue

            if allow_paragraph and raw_title and idx != 0:
                slides.append(self._make_paragraph_slide(t("Overview"), raw_title, notes))

        return PresentationPlan(title=presentation_title, slides=slides[:MAX_SLIDES])

    def _make_section_slide(self, title: str) -> SlideSpec:
        return SlideSpec(
            layout="section_slide",
            title=title,
            plugins=[],
        )

    def _make_title_slide(self, title: str, subtitle: str = "") -> SlideSpec:
        return SlideSpec(
            layout="title_slide",
            title=title,
            subtitle=subtitle or "Generated from prompt",
            plugins=[],
        )

    def _make_paragraph_slide(self, title: str, text: str, notes: str = "") -> SlideSpec:
        plugins: List[SlidePlugin] = [
            SlidePluginParagraph(
                type="paragraph",
                data={"title": title or "Overview", "text": text, "top": 1.45, "height": 3.8, "font_size": 18},
            )
        ]
        if notes:
            plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))
        return SlideSpec(layout="title_content", title=title or "Overview", plugins=plugins)

    def _make_bullets_slide(self, title: str, points: List[str], notes: str = "") -> SlideSpec:
        plugins: List[SlidePlugin] = [
            SlidePluginBullets(
                type="bullets",
                data={"title": title or "Key Points", "points": points, "top": 1.55, "height": 4.85},
            )
        ]
        if notes:
            plugins.append(SlidePluginNotes(type="notes", data={"notes": notes}))
        return SlideSpec(layout="bullets_slide", title=title or "Key Points", plugins=plugins)

    def _make_chart_slide(self, title: str, chart_payload: Dict[str, Any]) -> SlideSpec:
        return SlideSpec(layout="chart_slide", title=title or "Chart", plugins=[SlidePluginChart(type="chart", data=chart_payload)])

    def _make_image_slide(self, title: str, image_path: str) -> SlideSpec:
        return SlideSpec(
            layout="image_slide",
            title=title or "Visual",
            plugins=[SlidePluginImage(type="image", data={"path": image_path, "caption": title or "Visual", "title": title or "Visual"})],
        )

    def _make_table_slide(self, title: str, table_payload: Dict[str, Any]) -> SlideSpec:
        return SlideSpec(layout="table_slide", title=title or "Table", plugins=[SlidePluginTable(type="table", data=table_payload)])

    def _make_mixed_slide(self, title: str, plugins: List[SlidePlugin], notes: str = "") -> SlideSpec:
        if notes:
            plugins = plugins + [SlidePluginNotes(type="notes", data={"notes": notes})]
        return SlideSpec(layout="mixed_content_slide", title=title, plugins=plugins)

    def extract_structured_slides(self, prompt: str) -> List[str]:
        pattern = re.compile(
            r"(?:^|\n)\s*slide\s*\d+\s*[:\-]?\s*(.*?)(?=(?:\n\s*slide\s*\d+\s*[:\-]?)|$)",
            re.IGNORECASE | re.DOTALL,
        )
        return [block.strip() for block in pattern.findall(prompt) if block.strip()]

    def extract_overall_title(self, prompt: str, blocks: List[str]) -> str:
        if blocks:
            first_title = self.extract_section_value(blocks[0], "title")
            if first_title:
                return first_title
        m = re.search(r'presentation on\s*["“](.*?)["”]', prompt, re.IGNORECASE | re.DOTALL)
        if m:
            return normalize_whitespace(m.group(1))
        first_line = normalize_whitespace(prompt.split("\n", 1)[0])
        return first_line if len(first_line) <= 60 else first_line[:60].rstrip() + "..."

    def extract_section_value(self, text: str, key: str) -> Optional[str]:
        pattern = rf"(?im)^\s*{re.escape(key)}\b\s*[:\-]\s*(.+?)\s*$"
        m = re.search(pattern, text)
        if m:
            value = normalize_whitespace(m.group(1))
            return value or None
        return None

    def extract_heading_title(self, text: str) -> Optional[str]:
        for raw_line in text.splitlines():
            line = normalize_whitespace(raw_line)
            if not line:
                continue
            if re.match(r"^(title|subtitle|paragraph|bullets?|chart|chart type|values|categories|image|path|series name|table|notes|speaker notes|section)\b", line, re.IGNORECASE):
                continue
            if len(line) <= 50:
                return line
        return None

    def is_section_block(self, block: str, parsed: Dict[str, Any]) -> bool:
        if parsed.get("section_title"):
            return True

        lines = [normalize_whitespace(x) for x in block.splitlines() if normalize_whitespace(x)]
        if len(lines) != 1:
            return False

        line = lines[0]
        if len(line) <= 60 and line.upper() == line and any(ch.isalpha() for ch in line):
            if not parsed.get("title") and not parsed.get("subtitle"):
                return True
        return False

    def parse_slide_block(self, block: str) -> Dict[str, Any]:
        result: Dict[str, Any] = {
            "title": None,
            "subtitle": None,
            "section_title": None,
            "paragraph": "",
            "paragraph_lines": [],
            "bullets": [],
            "diagram": None,
            "chart_points": [],
            "chart_values": [],
            "chart_type": None,
            "chart_series": OrderedDict(),
            "chart_categories": [],
            "series_name": "Usage",
            "image_path": None,
            "table_headers": [],
            "table_rows": [],
            "notes_lines": [],
            "notes": "",
            "is_chart": False,
        }

        mode: Optional[str] = None
        current_series = "Usage"

        for raw_line in block.splitlines():
            line = raw_line.strip()
            if not line:
                continue

            m = re.match(r"^\s*title\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["title"] = clean_ai_instructions(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*subtitle\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["subtitle"] = clean_ai_instructions(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*section\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["section_title"] = clean_ai_instructions(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*(?:diagram|flowchart)\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["diagram"] = clean_ai_instructions(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*paragraph\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "paragraph"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["paragraph_lines"].append(tail)
                continue

            m = re.match(r"^\s*bullets?\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "bullets"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["bullets"].append(tail)
                continue

            m = re.match(r"^\s*table\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "table"
                tail = normalize_whitespace(m.group(1))
                if tail and "|" in tail:
                    row = self.parse_table_row(tail)
                    if row:
                        result["table_rows"].append(row)
                continue

            m = re.match(r"^\s*chart\s*type\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["chart_type"] = self.normalize_chart_type(m.group(1))
                result["is_chart"] = True
                mode = "chart"
                continue

            m = re.match(r"^\s*chart\b\s*[:\-]?\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                value = normalize_whitespace(m.group(1))
                if value and value.lower() in {"line", "bar", "column", "pie"}:
                    result["chart_type"] = self.normalize_chart_type(value)
                result["is_chart"] = True
                mode = "chart"
                continue

            m = re.match(r"^\s*series\s*name\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                result["series_name"] = current_series
                result["chart_series"].setdefault(current_series, OrderedDict())
                result["is_chart"] = True
                continue

            m = re.match(r"^\s*series\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                result["series_name"] = current_series
                result["chart_series"].setdefault(current_series, OrderedDict())
                result["is_chart"] = True
                continue

            m = re.match(r"^\s*(?:image|path)\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                result["image_path"] = normalize_whitespace(m.group(1))
                mode = None
                continue

            m = re.match(r"^\s*(?:notes|speaker\s*notes)\b\s*[:\-]?\s*(.*)$", line, re.IGNORECASE)
            if m:
                mode = "notes"
                tail = normalize_whitespace(m.group(1))
                if tail:
                    result["notes_lines"].append(tail)
                continue

            m = re.match(r"^\s*categories\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE)
            if m:
                raw = normalize_whitespace(m.group(1))
                cats = self.split_inline_list(raw)
                if cats:
                    result["chart_categories"] = cats
                continue

            if mode == "paragraph":
                if re.match(r"^(title|subtitle|bullets?|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section|diagram|flowchart)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    result["paragraph_lines"].append(line)
                    continue

            if mode == "bullets":
                bullet_match = re.match(r"^(?:[-*•]|\d+[.)])\s*(.*\S)$", line)
                if bullet_match:
                    point = clean_ai_instructions(bullet_match.group(1))
                    if point:
                        result["bullets"].append(point)
                    continue
                if not re.match(r"^(title|subtitle|paragraph|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section|diagram|flowchart)\b", line, re.IGNORECASE):
                    cleaned_pt = clean_ai_instructions(line)
                    if cleaned_pt:
                        result["bullets"].append(cleaned_pt)
                continue

            if mode == "table":
                if "|" in line:
                    row = self.parse_table_row(line)
                    if row:
                        result["table_rows"].append(row)
                    continue
                if re.match(r"^(title|subtitle|paragraph|bullets?|chart|chart type|image|path|notes|speaker notes|section|diagram|flowchart)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    if result["table_rows"]:
                        result["table_rows"][-1].append(normalize_whitespace(line))
                continue

            if mode == "chart":
                if m := re.match(r"^\s*type\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    result["chart_type"] = self.normalize_chart_type(m.group(1))
                    result["is_chart"] = True
                    continue
                if m := re.match(r"^\s*series\s*name\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    current_series = normalize_whitespace(m.group(1)) or f"Series {len(result['chart_series']) + 1}"
                    result["series_name"] = current_series
                    result["chart_series"].setdefault(current_series, OrderedDict())
                    result["is_chart"] = True
                    continue
                if m := re.match(r"^\s*categories\b\s*[:\-]\s*(.+?)\s*$", line, re.IGNORECASE):
                    cats = self.split_inline_list(normalize_whitespace(m.group(1)))
                    if cats:
                        result["chart_categories"] = cats
                    continue
                if m := re.match(r"^\s*([A-Za-z][A-Za-z0-9 _-]{0,30})\s*[:=]\s*(\d+(?:\.\d+)?)\s*$", line):
                    category = normalize_whitespace(m.group(1))
                    value = parse_number(m.group(2))
                    if category and value is not None:
                        result["chart_points"].append(category)
                        result["chart_values"].append(value)
                        result["chart_series"].setdefault(current_series, OrderedDict())
                        result["chart_series"][current_series][category] = value
                        result["is_chart"] = True
                    continue

            if mode == "notes":
                if re.match(r"^(title|subtitle|paragraph|bullets?|chart|chart type|categories|values|image|path|series name|table|section|diagram|flowchart)\b", line, re.IGNORECASE):
                    mode = None
                else:
                    result["notes_lines"].append(line)
                    continue

            if "|" in line:
                row = self.parse_table_row(line)
                if row and len(row) >= 2:
                    result["table_rows"].append(row)
                    continue

        if not result["title"]:
            raw_h = self.extract_heading_title(block)
            result["title"] = clean_ai_instructions(raw_h) if raw_h else None
        else:
            result["title"] = clean_ai_instructions(result["title"])

        if result["subtitle"]:
            result["subtitle"] = clean_ai_instructions(result["subtitle"])

        if result["paragraph_lines"]:
            result["paragraph"] = clean_ai_instructions(" ".join(result["paragraph_lines"]))

        if result["bullets"]:
            result["bullets"] = [clean_ai_instructions(b) for b in result["bullets"] if clean_ai_instructions(b)]

        if result["notes_lines"]:
            result["notes"] = clean_ai_instructions(" ".join(result["notes_lines"]))

        if result["chart_series"] and not result["chart_categories"]:
            seen = []
            for mapping in result["chart_series"].values():
                for cat in mapping.keys():
                    if cat not in seen:
                        seen.append(cat)
            result["chart_categories"] = seen
        if not result["chart_series"] and result["chart_points"]:
            series_name = result.get("series_name") or "Usage"
            result["chart_series"] = OrderedDict({series_name: OrderedDict(zip(result["chart_points"], result["chart_values"]))})
        if self.is_chart_block(block.lower()):
            result["is_chart"] = True
        return result

    def split_inline_list(self, raw: str) -> List[str]:
        parts = re.split(r"\s*[|,;/]\s*", normalize_whitespace(raw))
        return [p for p in (normalize_whitespace(x) for x in parts) if p]

    def parse_table_row(self, line: str) -> List[str]:
        parts = [normalize_whitespace(x) for x in line.strip().strip("|").split("|")]
        return [p for p in parts if p]

    def build_chart_payload(self, parsed: Dict[str, Any]) -> Dict[str, Any]:
        chart_series = parsed.get("chart_series") or OrderedDict()
        chart_type = parsed.get("chart_type") or "line"
        title = parsed.get("title") or "Growth Chart"
        categories = parsed.get("chart_categories") or []
        series_name = parsed.get("series_name") or "Usage"

        if chart_series:
            return {
                "chart_type": chart_type,
                "title": title,
                "categories": categories,
                "series_map": chart_series,
                "series_name": series_name,
            }

        categories = parsed.get("chart_points") or categories or ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
        values = parsed.get("chart_values") or [12, 18, 24, 31, 39, 48]
        return {
            "chart_type": chart_type,
            "title": title,
            "categories": categories,
            "values": values,
            "series_name": series_name,
        }

    def build_table_payload(self, parsed: Dict[str, Any], title: str) -> Dict[str, Any]:
        rows = parsed.get("table_rows") or []
        headers = parsed.get("table_headers") or []
        data_rows = rows

        if rows:
            if not headers and len(rows) >= 2:
                headers = rows[0]
                data_rows = rows[1:]
            elif not headers:
                headers = [f"Column {i + 1}" for i in range(len(rows[0]))]
                data_rows = rows

        return {"title": title or "Table", "headers": headers, "rows": data_rows}

    def is_chart_block(self, text_l: str) -> bool:
        if re.search(r"\b(chart|graph)\b", text_l):
            return True
        if re.search(r"\b[A-Za-z]{3,9}\s*[:=]\s*\d+(?:\.\d+)?\b", text_l):
            return True
        return False

    def normalize_chart_type(self, text: str) -> str:
        t = normalize_whitespace(text).lower()
        if "line" in t:
            return "line"
        if "bar" in t:
            return "bar"
        if "pie" in t:
            return "pie"
        return "column"

    def extract_bullets(self, text: str) -> List[str]:
        bullets: List[str] = []
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            if re.match(r"^(title|subtitle|paragraph|chart|chart type|categories|values|image|path|series name|table|notes|speaker notes|section)\b", line, re.IGNORECASE):
                continue
            bullet_match = re.match(r"^(?:[-*•]|\d+[.)])\s+(.*\S)$", line)
            if bullet_match:
                cleaned = normalize_whitespace(bullet_match.group(1))
                if cleaned:
                    bullets.append(cleaned)
        if bullets:
            return bullets[:MAX_BULLETS_PER_SLIDE]
        chunks = re.split(r"[.;]\s+|\n+", text)
        for chunk in chunks:
            chunk = normalize_whitespace(chunk)
            if not chunk:
                continue
            if re.match(r"^(title|subtitle|paragraph|chart|image|path|table|notes)\b", chunk, re.IGNORECASE):
                continue
            if len(chunk) >= 12:
                bullets.append(chunk)
        return bullets[:MAX_BULLETS_PER_SLIDE]


# ---------------------------------------------------------------------
# Plugins
# ---------------------------------------------------------------------

class BasePlugin:
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        pass

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        self.apply(slide, plan, theme_name=None)
        return current_y


class TextPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        text = normalize_whitespace(plan.get("text", "") or plan.get("subtitle", "") or plan.get("title", ""))
        if not text:
            return
        box_spec = as_box(plan, Box(0.8, 1.5, 11.7, 0.5))
        box = slide.shapes.add_textbox(Inches(box_spec.left), Inches(box_spec.top), Inches(box_spec.width), Inches(box_spec.height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        user_font = plan.get("font_size")
        font_size = int(user_font) if user_font and str(user_font).isdigit() else 18
        p.font.size = Pt(font_size)
        p.font.bold = True
        custom_color = plan.get("font_color") or plan.get("color")
        p.font.color.rgb = hex_to_rgb(custom_color) if custom_color else palette["accent"]

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        text = normalize_whitespace(plan.get("text", "") or plan.get("subtitle", "") or plan.get("title", ""))
        if not text:
            return current_y
        box = slide.shapes.add_textbox(Inches(left_margin), Inches(current_y), Inches(content_width), Inches(0.5))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        user_font = plan.get("font_size")
        font_size = int(user_font) if user_font and str(user_font).isdigit() else 18
        p.font.size = Pt(font_size)
        p.font.bold = True
        custom_color = plan.get("font_color") or plan.get("color")
        p.font.color.rgb = hex_to_rgb(custom_color) if custom_color else palette["accent"]
        return current_y + 0.65


class ParagraphPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        text = normalize_whitespace(plan.get("text", ""))
        if not text:
            return

        user_font = plan.get("font_size")
        font_size = int(user_font) if user_font and str(user_font).isdigit() else best_font_size_for_paragraph(text, base=14)

        default_height = 0.6 if len(text) < 120 else (0.8 if len(text) < 250 else 1.1)
        box_spec = as_box(plan, Box(0.8, 1.5, 5.6, default_height))

        box = slide.shapes.add_textbox(Inches(box_spec.left), Inches(box_spec.top), Inches(box_spec.width), Inches(box_spec.height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = text
        custom_color = plan.get("font_color") or plan.get("color")
        text_color = hex_to_rgb(custom_color) if custom_color else palette["text"]
        configure_text_frame(tf, font_size=font_size, color=text_color)

        alignment = str(plan.get("alignment", "left")).lower()
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
        if alignment in align_map:
            for p in tf.paragraphs:
                p.alignment = align_map[alignment]

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        text = normalize_whitespace(plan.get("text", ""))
        if not text:
            return current_y

        user_font = plan.get("font_size")
        font_size = int(user_font) if user_font and str(user_font).isdigit() else best_font_size_for_paragraph(text, base=14)

        default_height = 0.6 if len(text) < 120 else (0.8 if len(text) < 250 else 1.1)
        box_spec = as_box(plan, Box(left_margin, current_y, content_width, default_height))

        box = slide.shapes.add_textbox(Inches(box_spec.left), Inches(box_spec.top), Inches(box_spec.width), Inches(box_spec.height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = text
        custom_color = plan.get("font_color") or plan.get("color")
        text_color = hex_to_rgb(custom_color) if custom_color else palette["text"]
        configure_text_frame(tf, font_size=font_size, color=text_color)

        alignment = str(plan.get("alignment", "left")).lower()
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
        if alignment in align_map:
            for p in tf.paragraphs:
                p.alignment = align_map[alignment]

        return max(current_y, box_spec.top) + box_spec.height + 0.15


class BulletsPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        points = safe_list(plan.get("points"))
        if not points:
            return

        user_font = plan.get("font_size")
        bullet_font = int(user_font) if user_font and str(user_font).isdigit() else best_font_size_for_bullets(points, base=14)

        default_height = max(0.6, 0.28 * len(points))
        box_spec = as_box(plan, Box(0.8, 1.5, 5.6, default_height))

        box = slide.shapes.add_textbox(Inches(box_spec.left), Inches(box_spec.top), Inches(box_spec.width), Inches(box_spec.height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, point in enumerate(points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"•  {point}"
            p.level = 0
            p.space_after = Pt(2)

        custom_color = plan.get("font_color") or plan.get("color")
        text_color = hex_to_rgb(custom_color) if custom_color else palette["text"]
        configure_text_frame(tf, font_size=bullet_font, color=text_color)

        alignment = str(plan.get("alignment", "left")).lower()
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
        if alignment in align_map:
            for p in tf.paragraphs:
                p.alignment = align_map[alignment]

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        points = safe_list(plan.get("points"))
        if not points:
            return current_y

        user_font = plan.get("font_size")
        bullet_font = int(user_font) if user_font and str(user_font).isdigit() else best_font_size_for_bullets(points, base=14)

        default_height = max(0.6, 0.28 * len(points))
        box_spec = as_box(plan, Box(left_margin, current_y, content_width, default_height))

        box = slide.shapes.add_textbox(Inches(box_spec.left), Inches(box_spec.top), Inches(box_spec.width), Inches(box_spec.height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, point in enumerate(points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"•  {point}"
            p.level = 0
            p.space_after = Pt(2)

        custom_color = plan.get("font_color") or plan.get("color")
        text_color = hex_to_rgb(custom_color) if custom_color else palette["text"]
        configure_text_frame(tf, font_size=bullet_font, color=text_color)

        alignment = str(plan.get("alignment", "left")).lower()
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
        if alignment in align_map:
            for p in tf.paragraphs:
                p.alignment = align_map[alignment]

        return max(current_y, box_spec.top) + box_spec.height + 0.20


class ChartPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        chart_type = str(plan.get("chart_type", "column")).lower()
        categories = list(plan.get("categories", []))
        values = list(plan.get("values", []))
        series_map = plan.get("series_map") or {}
        series_name = normalize_whitespace(plan.get("series_name", "Usage"))
        top_pos = float(plan.get("top", 1.8))
        raw_box = as_box(plan, Box(0.9, top_pos, 8.5, 3.8))
        safe_top = min(raw_box.top, 4.5)
        safe_height = min(raw_box.height, round(6.8 - safe_top, 2))
        box = Box(raw_box.left, safe_top, raw_box.width, max(1.5, safe_height))

        chart_data = CategoryChartData()
        if series_map:
            if not categories:
                seen: List[str] = []
                for mapping in series_map.values():
                    if isinstance(mapping, dict):
                        for cat in mapping.keys():
                            if cat not in seen:
                                seen.append(cat)
                categories = seen
            if not categories:
                categories = ["Phase 1", "Phase 2", "Phase 3", "Phase 4"]
            chart_data.categories = categories
            for s_name, mapping in series_map.items():
                if isinstance(mapping, dict):
                    series_values = []
                    for cat in categories:
                        val = mapping.get(cat, 0)
                        num = parse_number(re.sub(r"[^\d.-]", "", str(val))) if not isinstance(val, (int, float)) else float(val)
                        series_values.append(num if num is not None else 0.0)
                    chart_data.add_series(str(s_name), series_values)
                elif isinstance(mapping, (list, tuple)):
                    series_values = []
                    for v in mapping[:len(categories)]:
                        num = parse_number(re.sub(r"[^\d.-]", "", str(v))) if not isinstance(v, (int, float)) else float(v)
                        series_values.append(num if num is not None else 0.0)
                    if len(series_values) < len(categories):
                        series_values += [0.0] * (len(categories) - len(series_values))
                    chart_data.add_series(str(s_name), series_values)
        else:
            clean_values = []
            for v in values:
                num = parse_number(re.sub(r"[^\d.-]", "", str(v))) if not isinstance(v, (int, float)) else float(v)
                clean_values.append(num if num is not None else 0.0)

            if not clean_values or all(v == 0.0 for v in clean_values):
                values = [round(25.0 * (i + 1) * 1.2, 1) for i in range(len(categories) or 4)]
                if not categories:
                    categories = ["Phase 1", "Phase 2", "Phase 3", "Phase 4"]
                if "[Illustrative Data]" not in series_name:
                    series_name = f"{series_name} [Illustrative Data]"
            else:
                values = clean_values
                if not categories:
                    categories = [f"Phase {i+1}" for i in range(len(values))]

            n = min(len(categories), len(values))
            categories = categories[:n]
            values = values[:n]

            chart_data.categories = categories
            chart_data.add_series(series_name, values)

        chart_kind = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type in {"line", "trend"}:
            chart_kind = XL_CHART_TYPE.LINE_MARKERS
        elif chart_type in {"bar", "bar_horizontal"}:
            chart_kind = XL_CHART_TYPE.BAR_CLUSTERED
        elif chart_type in {"pie"}:
            chart_kind = XL_CHART_TYPE.PIE
        elif chart_type in {"area"}:
            chart_kind = XL_CHART_TYPE.AREA
        elif chart_type in {"donut", "doughnut"}:
            chart_kind = XL_CHART_TYPE.DOUGHNUT
        elif chart_type in {"column"}:
            chart_kind = XL_CHART_TYPE.COLUMN_CLUSTERED

        try:
            chart_shape = slide.shapes.add_chart(chart_kind, Inches(box.left), Inches(box.top), Inches(box.width), Inches(box.height), chart_data)
            chart = chart_shape.chart

            title_text = plan.get("title")
            if title_text and plan.get("show_title", True):
                chart.has_title = True
                chart.chart_title.text_frame.text = str(title_text)

            show_legend = plan.get("show_legend", True)
            chart.has_legend = show_legend
            if show_legend and "legend_position" in plan:
                pos_key = str(plan["legend_position"]).lower()
                pos_map = {
                    "top": XL_LEGEND_POSITION.TOP,
                    "bottom": XL_LEGEND_POSITION.BOTTOM,
                    "left": XL_LEGEND_POSITION.LEFT,
                    "right": XL_LEGEND_POSITION.RIGHT,
                }
                if pos_key in pos_map:
                    chart.legend.position = pos_map[pos_key]

            show_data_labels = plan.get("show_data_labels", False)
            if show_data_labels and len(chart.plots) > 0:
                chart.plots[0].has_data_labels = True

        except Exception as exc:
            logger.warning("Failed to render PPT chart: %s", exc)
            fallback_box = slide.shapes.add_textbox(Inches(box.left), Inches(box.top), Inches(box.width), Inches(1.5))
            tf = fallback_box.text_frame
            tf.text = f"[{series_name} Chart]\nData: " + ", ".join(f"{c}: {v}" for c, v in zip(categories, values if not series_map else []))
            try:
                tf.paragraphs[0].runs[0].font.color.rgb = palette["text"]
            except Exception:
                pass

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        self.apply(slide, {**plan, "top": current_y, "box": {"left": left_margin, "top": current_y, "width": 8.5, "height": 3.8}}, theme_name=None)
        return current_y + 4.0


class DiagramPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        diagram_text = plan.get("diagram", "") or plan.get("text", "") or "Input ➔ Process ➔ Output"
        top_pos = float(plan.get("top", 1.5))
        box = as_box(plan, Box(0.8, top_pos, 11.7, 1.5))

        # Header label
        badge_box = slide.shapes.add_textbox(Inches(box.left), Inches(box.top), Inches(box.width), Inches(0.35))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        p_b.text = "🔄 PROCESS & WORKFLOW DIAGRAM"
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = palette["accent"]

        # Split into steps
        raw_steps = re.split(r"\s*(?:➔|->|-->|\|)\s*", diagram_text)
        steps = [clean_ai_instructions(s).strip("[] ") for s in raw_steps if clean_ai_instructions(s).strip("[] ")]

        if len(steps) >= 2 and len(steps) <= 5:
            card_top = box.top + 0.4
            card_height = 0.65
            total_width = box.width
            num_steps = len(steps)
            gap = 0.25
            card_width = max(1.5, (total_width - (gap * (num_steps - 1))) / num_steps)

            for i, step in enumerate(steps):
                c_left = box.left + i * (card_width + gap)
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c_left), Inches(card_top), Inches(card_width), Inches(card_height))
                try:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = palette["accent"]
                    shape.line.color.rgb = palette["text"]
                    shape.line.width = Pt(1)
                except Exception:
                    pass

                tf = shape.text_frame
                tf.word_wrap = True
                tf.clear()
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = step
                set_run_style(run, font_size=11, bold=True, color=palette["background"])

                if i < num_steps - 1:
                    arrow_box = slide.shapes.add_textbox(Inches(c_left + card_width), Inches(card_top + 0.15), Inches(gap), Inches(0.35))
                    ap = arrow_box.text_frame.paragraphs[0]
                    ap.alignment = PP_ALIGN.CENTER
                    ap.text = "➔"
                    ap.font.size = Pt(14)
                    ap.font.color.rgb = palette["accent"]
        else:
            diag_box = slide.shapes.add_textbox(Inches(box.left), Inches(box.top + 0.35), Inches(box.width), Inches(box.height - 0.35))
            tf = diag_box.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = diagram_text
            set_run_style(run, font_size=14, bold=True, color=palette["text"])

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        self.apply(slide, {**plan, "top": current_y, "box": {"left": left_margin, "top": current_y, "width": content_width, "height": 1.5}}, theme_name=None)
        return current_y + 1.6


class ImagePlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        url = normalize_whitespace(plan.get("url", ""))
        path = normalize_whitespace(plan.get("path", ""))
        caption = normalize_whitespace(plan.get("caption", ""))
        top_pos = float(plan.get("top", 1.8))
        raw_box = as_box(plan, Box(0.8, top_pos, 6.6, 3.5))

        safe_top = min(raw_box.top, 4.5)
        caption_space = 0.35
        safe_height = min(raw_box.height, max(1.5, round(6.5 - safe_top - caption_space, 2)))
        box = Box(raw_box.left, safe_top, raw_box.width, safe_height)

        target_source = url or path
        safe_path = sanitize_image_path(target_source)
        if not safe_path:
            query = caption or plan.get("title") or "presentation visual"
            fetched = fetch_unsplash_url(query) or fetch_unsplash_image_for_topic(query)
            if fetched:
                safe_path = sanitize_image_path(fetched)

        if safe_path:
            try:
                from PIL import Image as PILImage
                with PILImage.open(safe_path) as img:
                    img_w, img_h = img.size

                aspect = img_w / img_h if img_h > 0 else 1.0
                target_aspect = box.width / box.height if box.height > 0 else 1.0

                if aspect > target_aspect:
                    render_w = box.width
                    render_h = box.width / aspect
                else:
                    render_h = box.height
                    render_w = box.height * aspect

                pos_left = box.left + (box.width - render_w) / 2
                pos_top = box.top + (box.height - render_h) / 2

                slide.shapes.add_picture(
                    safe_path,
                    Inches(pos_left),
                    Inches(pos_top),
                    width=Inches(render_w),
                    height=Inches(render_h),
                )

                display_label = caption or plan.get("title") or "Visual"
                cap_top = min(6.8, pos_top + render_h + 0.05)
                cap_box = slide.shapes.add_textbox(Inches(pos_left), Inches(cap_top), Inches(render_w), Inches(0.35))
                cap_tf = cap_box.text_frame
                cap_tf.word_wrap = True
                p = cap_tf.paragraphs[0]
                p.text = f"📷 {display_label}"
                p.alignment = PP_ALIGN.CENTER
                set_run_style(p.runs[0] if p.runs else p.add_run(), font_size=11, bold=True, color=palette["accent"])
                return
            except Exception as exc:
                logger.warning("Failed to insert picture %s: %s", safe_path, exc)
                safe_path = None

        if not safe_path:
            box_shape = slide.shapes.add_textbox(Inches(box.left), Inches(box.top), Inches(box.width), Inches(box.height))
            tf = box_shape.text_frame
            tf.text = f"Visual: {caption or path or 'Topic'}"
            tf.paragraphs[0].font.size = Pt(18)
            try:
                tf.paragraphs[0].runs[0].font.color.rgb = palette["text"]
            except Exception:
                pass

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        avail_h = max(1.5, round(6.5 - current_y - 0.35, 2))
        img_h = min(3.4, avail_h)
        img_w = min(6.6, content_width)
        self.apply(slide, {**plan, "top": current_y, "box": {"left": left_margin, "top": current_y, "width": img_w, "height": img_h}}, theme_name=None)
        return current_y + img_h + 0.45


class TablePlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        palette = get_theme_palette(theme_name)
        headers = safe_list(plan.get("headers"))
        rows = safe_list(plan.get("rows"))
        top_pos = float(plan.get("top", 1.6))
        raw_box = as_box(plan, Box(0.8, top_pos, 11.7, 3.5))

        safe_top = min(raw_box.top, 4.2)
        safe_height = min(raw_box.height, max(1.5, round(6.6 - safe_top, 2)))
        box = Box(raw_box.left, safe_top, raw_box.width, safe_height)

        if not headers or not rows:
            box_shape = slide.shapes.add_textbox(Inches(1.0), Inches(safe_top), Inches(8.0), Inches(1.0))
            tf = box_shape.text_frame
            tf.text = "Table data not found or incomplete."
            tf.paragraphs[0].font.size = Pt(18)
            return

        custom_header_bg = plan.get("header_bg")
        custom_header_color = plan.get("header_color")
        custom_cell_bg = plan.get("cell_bg")
        custom_cell_color = plan.get("cell_color")
        
        raw_h_fs = plan.get("header_font_size")
        font_size_header = int(raw_h_fs) if raw_h_fs and str(raw_h_fs).isdigit() else 12

        raw_c_fs = plan.get("cell_font_size", plan.get("font_size"))
        font_size_cell = int(raw_c_fs) if raw_c_fs and str(raw_c_fs).isdigit() else 10

        align_opt = str(plan.get("align", plan.get("alignment", "left"))).lower()
        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
        cell_align = align_map.get(align_opt, PP_ALIGN.LEFT)

        # Resolve Header Fill & Text Colors
        if custom_header_bg:
            hdr_bg_rgb = hex_to_rgb(str(custom_header_bg))
        else:
            hdr_bg_rgb = palette.get("table_header_bg") or palette["accent"]

        if custom_header_color:
            hdr_txt_rgb = hex_to_rgb(str(custom_header_color))
        else:
            hdr_txt_rgb = palette.get("table_header_text") or (RGBColor(15, 23, 42) if is_light_color(hdr_bg_rgb) else RGBColor(255, 255, 255))

        # Resolve Row Fills (Alternating Zebra Striping)
        if custom_cell_bg:
            base_row_rgb1 = hex_to_rgb(str(custom_cell_bg))
            base_row_rgb2 = base_row_rgb1
        else:
            base_row_rgb1 = palette.get("table_row_bg1") or RGBColor(30, 41, 59)
            base_row_rgb2 = palette.get("table_row_bg2") or RGBColor(15, 23, 42)

        default_row_txt = palette.get("table_row_text") or palette["text"]

        cols = len(headers)
        row_count = len(rows) + 1
        table_shape = slide.shapes.add_table(row_count, cols, Inches(box.left), Inches(box.top), Inches(box.width), Inches(box.height))
        table = table_shape.table

        # Format Header Row
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(header)
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_bg_rgb
            except Exception:
                pass
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.alignment = cell_align
                for run in p.runs:
                    set_run_style(run, font_size=font_size_header, bold=True, color=hdr_txt_rgb)

        # Format Data Rows with automatic high contrast text
        for r, row in enumerate(rows, start=1):
            row_bg_rgb = base_row_rgb1 if r % 2 == 1 else base_row_rgb2
            if custom_cell_color:
                row_txt_rgb = hex_to_rgb(str(custom_cell_color))
            else:
                row_txt_rgb = default_row_txt if not is_light_color(row_bg_rgb) else (RGBColor(15, 23, 42) if is_light_color(palette["background"]) else RGBColor(255, 255, 255))

            if isinstance(row, (list, tuple)):
                row_list = list(row)
            elif hasattr(row, "__iter__") and not isinstance(row, (str, bytes, dict)):
                row_list = list(row)
            else:
                row_list = [str(row)]

            values = row_list + [""] * max(0, cols - len(row_list))
            values = values[:cols]
            for c, value in enumerate(values):
                cell = table.cell(r, c)
                cell.text = str(value)
                try:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_bg_rgb
                except Exception:
                    pass
                tf = cell.text_frame
                tf.word_wrap = True
                for p in tf.paragraphs:
                    p.alignment = cell_align
                    for run in p.runs:
                        set_run_style(run, font_size=font_size_cell, color=row_txt_rgb)

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        avail_h = max(1.5, round(6.5 - current_y, 2))
        tbl_h = min(3.4, avail_h)
        self.apply(slide, {**plan, "top": current_y, "box": {"left": left_margin, "top": current_y, "width": content_width, "height": tbl_h}}, theme_name=None)
        return current_y + tbl_h + 0.35


class NotesPlugin(BasePlugin):
    def apply(self, slide, plan: Dict[str, Any], theme_name: Optional[str] = None) -> None:
        set_slide_notes(slide, plan.get("notes", ""))

    def apply_with_y(
        self,
        slide,
        plan: Dict[str, Any],
        current_y: float,
        left_margin: float,
        content_width: float,
        palette: Dict[str, RGBColor],
    ) -> float:
        set_slide_notes(slide, plan.get("notes", ""))
        return current_y


PLUGIN_REGISTRY: Dict[str, BasePlugin] = {
    "text": TextPlugin(),
    "paragraph": ParagraphPlugin(),
    "bullets": BulletsPlugin(),
    "chart": ChartPlugin(),
    "image": ImagePlugin(),
    "table": TablePlugin(),
    "notes": NotesPlugin(),
    "diagram": DiagramPlugin(),
}


# ---------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------

class PptRenderer:
    def __init__(self, template_file: str = DEFAULT_TEMPLATE_FILE):
        self.template_file = template_file

    def auto_select_layout(self, slide_spec: SlideSpec, slide_index: int) -> str:
        plugin_types = {p.type for p in slide_spec.plugins}
        content_types = plugin_types - {"notes"}

        if slide_spec.layout:
            return slide_spec.layout

        if slide_index == 0 and content_types == {"text"}:
            return "title_slide"

        mixed_kinds = {"paragraph", "bullets", "image", "chart", "table"}
        if len(content_types & mixed_kinds) >= 2:
            return "mixed_content_slide"

        if "chart" in content_types:
            return "chart_slide"
        if "image" in content_types:
            return "image_slide"
        if "table" in content_types:
            return "table_slide"
        if content_types == {"bullets"}:
            return "bullets_slide"
        if "paragraph" in content_types:
            return "title_content"
        if "text" in content_types:
            return "title_content"
        return "title_content"

    def render(self, plan: PresentationPlan, content_theme: Optional[str] = None, visual_style: Optional[str] = None) -> Presentation:
        prs = ensure_template_prs(self.template_file)
        layout_registry = get_layout_registry(self.template_file)
        active_theme = plan.theme or content_theme

        for idx, slide_spec in enumerate(plan.slides):
            layout_key = self.auto_select_layout(slide_spec, idx)
            layout_spec = layout_registry.get(layout_key, LayoutSpec(layout_index=0))
            layout_index = layout_spec.layout_index
            if layout_index >= len(prs.slide_layouts):
                layout_index = 0

            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)

            # Remove built-in template placeholders so "Click to add title" doesn't overlap!
            for sp in list(slide.placeholders):
                try:
                    sp._element.getparent().remove(sp._element)
                except Exception:
                    pass

            apply_background_theme(slide, active_theme, visual_style=visual_style)

            palette = get_theme_palette(active_theme)
            is_cover = idx == 0 or slide_spec.layout in {"title_slide", "section_slide"}
            current_y = 2.0 if is_cover else 0.35
            slide_width_in = float(prs.slide_width / Inches(1))
            left_margin = 0.6
            content_width = max(6.0, slide_width_in - (left_margin * 2.0))

            # 1. Slide Badge ("SLIDE X OF Y") matching Theme & Contrast
            bg_is_light = is_light_color(palette["background"])
            badge_color = palette.get("badge") or palette["accent"]
            if bg_is_light and is_light_color(badge_color):
                badge_color = palette["text"]

            badge_box = slide.shapes.add_textbox(Inches(left_margin), Inches(current_y), Inches(content_width), Inches(0.25))
            p_b = badge_box.text_frame.paragraphs[0]
            p_b.text = f"SLIDE {idx + 1} OF {len(plan.slides)}"
            p_b.font.size = Pt(9)
            p_b.font.bold = True
            p_b.font.color.rgb = badge_color
            if is_cover:
                p_b.alignment = PP_ALIGN.CENTER
            current_y += 0.35

            # 2. Main Title Rendering
            title_text = slide_spec.title or (plan.title if idx == 0 else "")
            if title_text:
                default_title_size = 32 if is_cover else 22
                title_font_size = slide_spec.title_font_size or default_title_size
                title_color = hex_to_rgb(slide_spec.title_color) if slide_spec.title_color else palette["text"]
                title_bold = slide_spec.title_bold if slide_spec.title_bold is not None else True
                align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
                default_align = "center" if is_cover else "left"
                t_align = align_map.get(str(slide_spec.title_align or default_align).lower(), PP_ALIGN.CENTER if is_cover else PP_ALIGN.LEFT)

                t_box = slide.shapes.add_textbox(Inches(left_margin), Inches(current_y), Inches(content_width), Inches(0.70))
                tf_t = t_box.text_frame
                tf_t.word_wrap = True
                p_t = tf_t.paragraphs[0]
                p_t.text = title_text
                p_t.alignment = t_align
                set_run_style(p_t.runs[0] if p_t.runs else p_t.add_run(), font_size=title_font_size, bold=title_bold, color=title_color)
                current_y += 0.75

            # 3. Subtitle Rendering
            if slide_spec.subtitle:
                sub_font_size = slide_spec.subtitle_font_size or 15
                sub_color = hex_to_rgb(slide_spec.subtitle_color) if slide_spec.subtitle_color else RGBColor(148, 163, 184)
                align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT}
                default_align = "center" if is_cover else "left"
                s_align = align_map.get(str(slide_spec.subtitle_align or default_align).lower(), PP_ALIGN.CENTER if is_cover else PP_ALIGN.LEFT)

                sub_box = slide.shapes.add_textbox(Inches(left_margin), Inches(current_y), Inches(content_width), Inches(0.45))
                tf_s = sub_box.text_frame
                tf_s.word_wrap = True
                p_s = tf_s.paragraphs[0]
                p_s.text = slide_spec.subtitle
                p_s.alignment = s_align
                set_run_style(p_s.runs[0] if p_s.runs else p_s.add_run(), font_size=sub_font_size, bold=False, color=sub_color)
                current_y += 0.50

            current_y += 0.05 # Padding gap

            # 4. Plugins rendering
            for plugin in slide_spec.plugins:
                handler = PLUGIN_REGISTRY.get(plugin.type)
                if handler is None:
                    continue
                if (slide_spec.layout == "mixed_content_slide" or len(slide_spec.plugins) >= 2) and "box" in plugin.data:
                    handler.apply(slide, plugin.data, theme_name=active_theme)
                else:
                    next_y = handler.apply_with_y(slide, plugin.data, current_y=current_y, left_margin=left_margin, content_width=content_width, palette=palette)
                    if next_y is not None and next_y > current_y:
                        current_y = next_y

        return prs


# ---------------------------------------------------------------------
# Storage
# ---------------------------------------------------------------------

def save_presentation(prs: Presentation, title: str) -> str:
    file_id = uuid.uuid4().hex
    filename = f"{safe_filename(title)}_{file_id}.pptx"
    file_path = OUTPUT_DIR / filename
    prs.save(str(file_path))
    return str(file_path)


# ---------------------------------------------------------------------
# Service
# ---------------------------------------------------------------------

def ensure_plan_images(plan: PresentationPlan, allow_image: bool = True) -> PresentationPlan:
    """Auto-populate image URLs for any image plugins in the plan, and enrich 2-3 suitable slides with HD Unsplash images if images are allowed."""
    image_count = 0
    for slide in plan.slides:
        for plugin in slide.plugins:
            if plugin.type == "image":
                image_count += 1
                data = dict(plugin.data)
                url = data.get("url") or data.get("path") or ""
                caption = data.get("caption") or data.get("title") or slide.title or "Visual"

                if not url or (not url.startswith("http") and not Path(url).exists()):
                    query = f"{plan.title} {caption}"
                    live_url = fetch_unsplash_url(query) or fetch_unsplash_url(caption)
                    local_path = fetch_unsplash_image(query) or fetch_unsplash_image(caption)
                    url = live_url or local_path or url

                data["url"] = url
                data["path"] = url
                plugin.data = data

    # Enrich suitable slides with images if fewer than 2 images exist
    if allow_image and image_count < 2:
        for idx, slide in enumerate(plan.slides):
            if idx == 0:  # Skip title cover slide
                continue
            if image_count >= 3:
                break
            
            plugin_types = {p.type for p in slide.plugins}
            # Attach image to text/paragraph/bullet slides that do not already have chart/table/diagram
            if "image" not in plugin_types and not (plugin_types & {"chart", "table", "diagram"}):
                query = f"{plan.title} {slide.title or 'technology'}"
                live_url = fetch_unsplash_url(query) or fetch_unsplash_url(slide.title or "innovation")
                if live_url:
                    img_plugin = SlidePluginImage(
                        type="image",
                        data={"url": live_url, "path": live_url, "caption": slide.title or "Visual Highlight", "title": slide.title or "Visual Highlight"}
                    )
                    slide.plugins.append(img_plugin)
                    if slide.layout == "title_content":
                        slide.layout = "mixed_content_slide"
                    image_count += 1
                    
                    # Update box geometry for 2-column side-by-side layout
                    box_list = MixedLayoutResolver.resolve_list([p.type for p in slide.plugins])
                    for plugin, b in zip(slide.plugins, box_list):
                        if b:
                            plugin.data["box"] = {"left": b.left, "top": b.top, "width": b.width, "height": b.height}

    return plan


class PresentationService:
    def __init__(self) -> None:
        self.planner = PromptPlanner()

    def generate(self, req: GenerateRequest) -> tuple[str, PresentationPlan, str]:
        template_file = resolve_template_path(req.template_name)
        # Renderer carries the chosen template. Keep it request-local so two
        # simultaneous users cannot accidentally render with each other's
        # template.
        renderer = PptRenderer(template_file=template_file)

        if req.plan is not None:
            if isinstance(req.plan, dict):
                plan = PresentationPlan(**req.plan)
            else:
                plan = req.plan
        else:
            # Gemini produces researched, slide-by-slide content. The local prompt
            # planner is still used to validate/normalize that content into plugins.
            planning_prompt = build_gemini_slide_script(req) or req.prompt

            plan = self.planner.plan(
                planning_prompt,
                include_title_slide=req.include_title_slide,
                allow_bullets=req.allow_bullets,
                allow_paragraph=req.allow_paragraph,
                allow_chart=req.allow_chart,
                allow_image=req.allow_image,
                allow_section_slide=req.allow_section_slide,
                allow_table=req.allow_table,
                smart_mode=req.smart_mode,
                slide_types=normalize_slide_types(req.slide_types),
                target_slide_count=req.slide_count,
            )

        plan = ensure_plan_images(plan, allow_image=req.allow_image)

        content_theme = normalize_whitespace(req.content_theme or req.background_theme or "")
        if not content_theme or content_theme.lower() in {"auto", "detect"}:
            content_theme = detect_theme(req.prompt)

        visual_style = normalize_whitespace(req.visual_style or "")
        if not visual_style or visual_style.lower() in {"auto", "detect"}:
            visual_style = detect_visual_style(req.prompt)

        prs = renderer.render(plan, content_theme=content_theme, visual_style=visual_style)
        file_path = save_presentation(prs, plan.title)
        return file_path, plan, plan.title


service = PresentationService()


# ---------------------------------------------------------------------
# API
# ---------------------------------------------------------------------

@router.get("/health")
def health() -> Dict[str, str]:
    return {"status": "ok"}


@router.get("/unsplash/search")
def search_unsplash_image(query: str) -> Dict[str, str]:
    """Search Unsplash for an exact topic query and return live HD image URL."""
    url = fetch_unsplash_url(query) or fetch_unsplash_image(query) or ""
    return {"query": query, "url": url}


@router.post("/plan", response_model=PresentationPlan)
async def preview_plan(req: GenerateRequest) -> PresentationPlan:
    try:
        planning_prompt = await run_in_threadpool(build_gemini_slide_script, req)
        plan = await run_in_threadpool(
            service.planner.plan,
            planning_prompt or req.prompt,
            include_title_slide=req.include_title_slide,
            allow_bullets=req.allow_bullets,
            allow_paragraph=req.allow_paragraph,
            allow_chart=req.allow_chart,
            allow_image=req.allow_image,
            allow_section_slide=req.allow_section_slide,
            allow_table=req.allow_table,
            smart_mode=req.smart_mode,
            slide_types=normalize_slide_types(req.slide_types),
            target_slide_count=req.slide_count,
        )
        return ensure_plan_images(plan, allow_image=req.allow_image)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc


@router.post("/generate", response_model=GenerateResponse)
async def generate_presentation(req: GenerateRequest, request: Request) -> GenerateResponse:
    try:
        file_path, _plan, _title = await run_in_threadpool(service.generate, req)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    job_id = uuid.uuid4().hex
    filename = Path(file_path).name
    download_url = str(request.url_for("download_ppt", file_name=filename))

    return GenerateResponse(job_id=job_id, status="completed", file_name=filename, download_url=download_url)


@router.post("/save", response_model=SaveResponse)
async def save_presentation_endpoint(req: GenerateRequest, request: Request) -> SaveResponse:
    """Save presentation to backend, process PPTX generation, and return saved presentation details."""
    try:
        file_path, _plan, _title = await run_in_threadpool(service.generate, req)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=str(exc)) from exc

    presentation_id = f"pres_{uuid.uuid4().hex[:12]}"
    filename = Path(file_path).name
    download_url = str(request.url_for("download_ppt", file_name=filename))

    return SaveResponse(
        presentation_id=presentation_id,
        status="saved",
        file_name=filename,
        download_url=download_url,
        message="Presentation saved successfully",
    )


@router.post("/refine-slide", response_model=RefineSlideResponse)
async def refine_slide_text(req: RefineSlideRequest) -> RefineSlideResponse:
    """Refine or polish slide text using AI."""
    if not req.text or not req.text.strip():
        return RefineSlideResponse(refined_text="")

    prompt = f"Refine and polish the following presentation text to be executive, clear, concise, and impact-driven:\n\n{req.text.strip()}"
    try:
        refined = await run_in_threadpool(generate_response, prompt)
        cleaned = (refined or "").strip()
        if cleaned.startswith("Gemini API key is not configured") or cleaned.startswith("Gemini error"):
            logger.warning("Refine slide AI failed: %s", cleaned)
            return RefineSlideResponse(refined_text=req.text)
        if cleaned.startswith('"') and cleaned.endswith('"'):
            cleaned = cleaned[1:-1].strip()
        return RefineSlideResponse(refined_text=cleaned or req.text)
    except Exception as exc:
        logger.warning("Refine slide AI call failed: %s", exc)
        return RefineSlideResponse(refined_text=req.text)


@router.get("/download/{file_name}", name="download_ppt")
def download_ppt(file_name: str) -> FileResponse:
    # Restrict downloads to files we generated; this also prevents path
    # traversal attempts through encoded path separators.
    if Path(file_name).name != file_name or not file_name.lower().endswith(".pptx"):
        raise HTTPException(status_code=404, detail="File not found")
    file_path = OUTPUT_DIR / file_name
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(
        path=str(file_path),
        filename=file_name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.get("/")
def root():
    return {
        "name": APP_NAME,
        "status": "ok",
        "endpoints": ["/plan", "/generate", "/download/{file_name}"],
        "max_slides": MAX_SLIDES,
    }
